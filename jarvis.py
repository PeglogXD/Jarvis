# ============================================================
# J.A.R.V.I.S. - Just A Rather Very Intelligent System
# ------------------------------------------------------------
# Archivo : jarvis.py
# Versión : 2.0.0
# Estado  : Producción — Upgrade masivo de funciones y UI
#
# Asistente personal avanzado con IA conversacional,
# automatización del PC, creación de documentos,
# herramientas de productividad y UI de nueva generación.
# ============================================================
import tkinter as tk
from tkinter import filedialog as fd
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPTpt
import openpyxl
import threading
import openai
from google import genai
import speech_recognition as sr
import os
import subprocess
import re
import pygame
import edge_tts
import asyncio
import webbrowser
import math
import time
import wave
import struct
import uuid
import unicodedata
import urllib.parse
import pyautogui
import requests
import io
import hashlib
import secrets
import string
import random
import tempfile
from core.logger import info, warning, error, exception
from config.settings import VERSION
from datetime import datetime, timedelta
from googleapiclient.discovery import build
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.options import Options
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.service import Service
from pushbullet import Pushbullet
import json
# (groq eliminado — se usa openai + google-generativeai)

# ─── INTENTO DE IMPORTAR MODULOS OPCIONALES ──────────────────────────────────
try:
    import psutil as _psutil
except ImportError:
    _psutil = None

# Ocultar ventanas CMD de procesos hijos (Ollama, taskkill, etc)
import sys
if sys.platform == "win32":
    import ctypes
    try:
        ctypes.windll.kernel32.SetConsoleWindowInfo(
            ctypes.windll.kernel32.GetConsoleWindow(), True,
            ctypes.byref((ctypes.c_short * 4)(0, 0, 0, 0))
        )
    except Exception:
        pass

NO_WINDOW = subprocess.CREATE_NO_WINDOW

# ─── CLAVES API (cargadas desde .env) ────────────────────────────────────────
from dotenv import load_dotenv
load_dotenv("C:/Jarvis/.env")

PB_API_KEY = os.getenv("PUSHBULLET_API_KEY")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
NVIDIA_API_KEY = os.getenv("NVIDIA_API_KEY")
YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY")

if not GEMINI_API_KEY and not NVIDIA_API_KEY:
    raise RuntimeError(
        "Falta GEMINI_API_KEY o NVIDIA_API_KEY. "
        "Crea el archivo C:/Jarvis/.env con al menos una de las claves."
    )

# ─── NVIDIA NIM (nemotron-3.5-lightning, compatible con OpenAI SDK) ───
nvidia_client = None
NVIDIA_MODEL = "nvidia/nemotron-3.5-lightning-30b-a3b"
if NVIDIA_API_KEY and NVIDIA_API_KEY.startswith("nvapi-"):
    nvidia_client = openai.OpenAI(
        api_key=NVIDIA_API_KEY,
        base_url="https://integrate.api.nvidia.com/v1"
    )

# ─── Gemini (Google AI) — PRIMARIO ───
gemini_client = None
GEMINI_MODEL = "gemini-3.5-flash"
if GEMINI_API_KEY:
    gemini_client = genai.Client(api_key=GEMINI_API_KEY)

reconocedor = sr.Recognizer()
mic_lock = threading.Lock()
memoria_path = "C:/Jarvis/memoria.txt"

# Ruta donde se guardan los apodos de salidas de audio ("Auriculares" -> dispositivo real)
SALIDAS_AUDIO_PATH = "C:/Jarvis/salidas_audio.json"

# Inicializar mezclador de audio
pygame.mixer.init()

hablando = False
mic_activo = True  # Micrófono activo al abrir (listener.py maneja el wake word)
modo_texto = False
voz_jarvis = "es-MX-JorgeNeural"  # fallback cuando OpenAI TTS no está disponible

# ─── NVIDIA SKILLS (knowledge base para consultas GPU/aceleradas) ───
NVIDIA_SKILLS_DIR = "C:/Jarvis/nvidia_skills"
NVIDIA_SKILL_KEYWORDS = {
    "cudf": ["cudf", "gpu dataframe", "rapids", "pandas gpu", "acelerar pandas", "gpu data"],
    "cuda": ["cuda", "gpu programming", "kernel gpu", "nvidia cuda"],
    "cuopt": ["cuopt", "optimizacion gpu", "vehicle routing", "linear programming gpu"],
    "rag": ["rag", "retrieval augmented", "vector database", "milvus", "elasticsearch", "nvidia rag", "deploy rag", "rag blueprint", "ingestor", "nim-llm", "ngc"],
}

def _detectar_nvidia_skill(prompt):
    """Detecta si el prompt requiere una skill de NVIDIA y retorna su contenido."""
    prompt_lower = prompt.lower()
    for skill_name, keywords in NVIDIA_SKILL_KEYWORDS.items():
        for kw in keywords:
            if kw in prompt_lower:
                skill_path = os.path.join(NVIDIA_SKILLS_DIR, f"{skill_name}.md")
                if os.path.exists(skill_path):
                    with open(skill_path, "r", encoding="utf-8") as f:
                        return f.read()
    return None
archivo_cargado = {"contenido": None, "nombre": None, "ruta_original": None}
estado_accion = "escuchando"  # controla color/animación de la bolita


# ─── LOG DE ERRORES (definido antes de usarse) ──────────────────────────────
def _log_error(origen, e):
    try:
        with open("C:/Jarvis/errores.log", "a", encoding="utf-8") as f:
            f.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {origen}: {e}\n")
    except Exception:
        pass


# ─── LIMPIEZA DE ARCHIVOS TEMPORALES DE VOZ ─────────────────────────────────
def limpiar_temporales():
    """Elimina archivos de voz residuales generados en sesiones previas"""
    try:
        directorio = "C:/Jarvis"
        if os.path.exists(directorio):
            for arch in os.listdir(directorio):
                if arch.startswith("voz_") and arch.endswith(".mp3"):
                    try:
                        os.remove(os.path.join(directorio, arch))
                    except Exception:
                        pass
    except Exception as e:
        _log_error("limpiar_temporales()", e)


# ─── GENERADOR DE EFECTOS DE SONIDO PROCEDIMENTALES (SONIDOS X10) ───────────
def generar_sfx_predeterminados():
    directorio_sonidos = "C:/Jarvis/Sounds"
    os.makedirs(directorio_sonidos, exist_ok=True)
    limpiar_temporales() 
    
    # 1. CLICK.WAV
    ruta_click = f"{directorio_sonidos}/click.wav"
    if not os.path.exists(ruta_click):
        try:
            with wave.open(ruta_click, 'wb') as w:
                w.setnchannels(1); w.setsampwidth(2); w.setframerate(22050)
                for i in range(350):
                    t = i / 22050
                    val = math.sin(2 * math.pi * 1900 * t) * math.exp(-320 * t)
                    w.writeframesraw(struct.pack('<h', int(val * 16000)))
        except Exception: pass

    # 2. LISTENING.WAV
    ruta_listening = f"{directorio_sonidos}/listening.wav"
    if not os.path.exists(ruta_listening):
        try:
            with wave.open(ruta_listening, 'wb') as w:
                w.setnchannels(1); w.setsampwidth(2); w.setframerate(22050)
                dur = 0.07
                n = int(dur * 22050)
                for i in range(n):
                    t = i / 22050
                    f = 950 + (1300 - 950) * (t / dur)
                    env = 1.0 - (i / n) if i > n * 0.75 else (i / (n * 0.25))
                    val = math.sin(2 * math.pi * f * t) * env
                    w.writeframesraw(struct.pack('<h', int(val * 13000)))
                for i in range(n):
                    t = i / 22050
                    f = 1400 + (1900 - 1400) * (t / dur)
                    env = 1.0 - (i / n) if i > n * 0.75 else (i / (n * 0.25))
                    val = math.sin(2 * math.pi * f * t) * env
                    w.writeframesraw(struct.pack('<h', int(val * 13000)))
        except Exception: pass

    # 3. SUCCESS.WAV
    ruta_success = f"{directorio_sonidos}/success.wav"
    if not os.path.exists(ruta_success):
        try:
            with wave.open(ruta_success, 'wb') as w:
                w.setnchannels(1); w.setsampwidth(2); w.setframerate(22050)
                freqs = [523.25, 659.25, 784.00, 1046.50, 1318.51]
                total_dur = 0.65
                num_frames = int(total_dur * 22050)
                for i in range(num_frames):
                    t = i / 22050
                    val = 0
                    for idx, f in enumerate(freqs):
                        start_t = idx * 0.045
                        if t >= start_t:
                            dt = t - start_t
                            val += math.sin(2 * math.pi * f * dt) * math.exp(-7.5 * dt)
                    val = max(-1.0, min(1.0, val / 2.2))
                    env = (num_frames - i) / (num_frames * 0.25) if i > num_frames * 0.75 else 1.0
                    w.writeframesraw(struct.pack('<h', int(val * 15000 * env)))
        except Exception: pass

    # 4. ERROR.WAV
    ruta_error = f"{directorio_sonidos}/error.wav"
    if not os.path.exists(ruta_error):
        try:
            with wave.open(ruta_error, 'wb') as w:
                w.setnchannels(1); w.setsampwidth(2); w.setframerate(22050)
                dur = 0.4
                n = int(dur * 22050)
                for i in range(n):
                    t = i / 22050
                    val = (math.sin(2 * math.pi * 140 * t) + math.sin(2 * math.pi * 144 * t)) / 2.0
                    if val > 0.35: val = 0.35
                    elif val < -0.35: val = -0.35
                    env = (n - i) / (n * 0.25) if i > n * 0.75 else (i / (n * 0.1))
                    w.writeframesraw(struct.pack('<h', int(val * 20000 * env)))
        except Exception: pass

    # 5. PROCESSING.WAV
    ruta_processing = f"{directorio_sonidos}/processing.wav"
    if not os.path.exists(ruta_processing):
        try:
            with wave.open(ruta_processing, 'wb') as w:
                w.setnchannels(1); w.setsampwidth(2); w.setframerate(22050)
                dur = 0.15
                n = int(dur * 22050)
                for i in range(n):
                    t = i / 22050
                    f = 650 - (250 * (t / dur))
                    env = math.exp(-14 * t)
                    val = math.sin(2 * math.pi * f * t) * env
                    w.writeframesraw(struct.pack('<h', int(val * 11000)))
        except Exception: pass

generar_sfx_predeterminados()

def reproducir_sfx(nombre):
    try:
        ruta = f"C:/Jarvis/Sounds/{nombre}.wav"
        if os.path.exists(ruta):
            snd = pygame.mixer.Sound(ruta)
            snd.set_volume(0.55)
            snd.play()
    except Exception as e:
        _log_error("reproducir_sfx()", e)


def preguntar_ia(prompt, max_tokens=1024):
    """Llamada síncrona a IA. Orden: Gemini → NVIDIA NIM → OpenAI."""
    # ─── Detectar skill NVIDIA ───
    nvidia_skill = _detectar_nvidia_skill(prompt)
    if nvidia_skill:
        prompt = prompt + "\n\n---\nCONOCIMIENTO NVIDIA CARGADO:\n" + nvidia_skill
        max_tokens = max(max_tokens, 2048)
    # 1. Gemini (Google AI — principal)
    if gemini_client:
        try:
            respuesta = gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
                config=genai.types.GenerateContentConfig(max_output_tokens=max_tokens)
            )
            return respuesta.text
        except Exception as e:
            _log_error("preguntar_ia() - Gemini", e)
    # 2. NVIDIA NIM (sin espera, directo)
    if nvidia_client:
        try:
            respuesta = nvidia_client.chat.completions.create(
                model=NVIDIA_MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=max_tokens
            )
            return respuesta.choices[0].message.content
        except Exception as e:
            _log_error("preguntar_ia() - NVIDIA NIM", e)
    return "⚠️ No hay proveedores de IA configurados o disponibles."

def leer_memoria():
    if os.path.exists(memoria_path):
        with open(memoria_path, "r", encoding="utf-8") as f:
            return f.read()
    return "Sin registros."

def escuchar():
    if not mic_activo: return ""
    if not mic_lock.acquire(timeout=1): return ""
    try:
        with sr.Microphone() as source:
            reconocedor.adjust_for_ambient_noise(source, duration=0.6)
            audio = reconocedor.listen(source, timeout=10, phrase_time_limit=15)
        texto = reconocedor.recognize_google(audio, language="es-ES")
        return texto
    except sr.WaitTimeoutError: return ""
    except sr.UnknownValueError: return ""
    except Exception as e:
        _log_error("escuchar()", e)
        return ""
    finally:
        mic_lock.release()

def _normalizar_texto(texto):
    """Quita acentos y pasa a minúsculas para que las palabras clave
    ('guardar', 'editar', etc.) se detecten aunque el reconocimiento de voz
    las devuelva con mayúsculas, tildes o pequeñas variaciones."""
    texto = texto.strip().lower()
    texto = "".join(
        c for c in unicodedata.normalize("NFD", texto)
        if unicodedata.category(c) != "Mn"
    )
    return texto

def _contiene_alguna(texto, palabras):
    return any(p in texto for p in palabras)

def escuchar_interrupcion():
    if not mic_lock.acquire(timeout=0.5): return
    try:
        r = sr.Recognizer()
        r.energy_threshold = 1400
        r.dynamic_energy_threshold = False
        with sr.Microphone() as source:
            time.sleep(0.5)
            while hablando:
                try:
                    audio = r.listen(source, timeout=0.5, phrase_time_limit=2)
                    interrumpir()
                    return
                except sr.WaitTimeoutError: pass
                except Exception as e: _log_error("escuchar_interrupcion()", e)
    except Exception as e: _log_error("escuchar_interrupcion()", e)
    finally: mic_lock.release()


def hablar(texto):
    """Genera voz con OpenAI TTS (natural) o edge-tts como fallback."""
    global hablando, voz_jarvis
    def _hablar():
        global hablando
        hablando = True
        identificador = uuid.uuid4().hex
        archivo = f"C:/Jarvis/voz_{identificador}.mp3"
        try:
            asyncio.run(edge_tts.Communicate(texto, voice=voz_jarvis).save(archivo))
            pygame.mixer.music.load(archivo)
            pygame.mixer.music.play()
            threading.Thread(target=escuchar_interrupcion, daemon=True).start()
            while pygame.mixer.music.get_busy() and hablando:
                time.sleep(0.1)
        except Exception as e:
            _log_error("hablar() thread", e)
        finally:
            pygame.mixer.music.unload()
            try:
                if os.path.exists(archivo): os.remove(archivo)
            except Exception: pass
            hablando = False

    threading.Thread(target=_hablar, daemon=False).start()
    t0 = time.time()
    while not hablando and (time.time() - t0 < 2.0): time.sleep(0.05)
    while hablando: time.sleep(0.05)


def interrumpir():
    global hablando
    pygame.mixer.music.stop()
    hablando = False

def leer_archivo(ruta):
    extension = ruta.lower().split(".")[-1]
    try:
        if extension in ["txt", "py", "js", "json", "html", "css", "md", "csv"]:
            with open(ruta, "r", encoding="utf-8") as f: return f.read()
        elif extension == "pdf":
            import fitz
            doc = fitz.open(ruta)
            return "\n".join([page.get_text() for page in doc])
        elif extension == "docx":
            from docx import Document
            doc = Document(ruta)
            return "\n".join([p.text for p in doc.paragraphs])
        elif extension in ["xlsx", "xls"]:
            import openpyxl
            wb = openpyxl.load_workbook(ruta)
            texto = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    texto += " | ".join([str(c) for c in row if c]) + "\n"
            return texto
        elif extension in ["png", "jpg", "jpeg", "webp"]:
            return f"[Imagen cargada: {ruta}]"
        else: return None
    except: return None


def cargar_archivo():
    global archivo_cargado
    from tkinter import filedialog
    ruta = filedialog.askopenfilename(
        title="Seleccionar archivo",
        filetypes=[
            ("Archivos de código", "*.py;*.js;*.json;*.html;*.css;*.md"),
            ("Documentos", "*.txt;*.pdf;*.docx;*.xlsx"),
            ("Todos los archivos", "*.*")
        ]
    )
    if not ruta: return
    reproducir_sfx("click")
    contenido = leer_archivo(ruta)
    nombre = ruta.replace("\\", "/").split("/")[-1]
    if contenido:
        archivo_cargado["contenido"] = contenido
        archivo_cargado["nombre"] = nombre
        archivo_cargado["ruta_original"] = ruta
        extension = nombre.lower().split(".")[-1]
        editables = ["txt", "py", "js", "json", "html", "css", "md", "csv"]
        if modo_texto and extension in editables:
            # En modo texto: abrir editor inline para archivos editables
            abrir_editor(ruta, nombre, contenido)
            mostrar_toast(f"Editor abierto: {nombre}", "info")
        else:
            # Modo voz o archivo no editable: mostrar en chat
            archivo_label.config(text=f"📎 {nombre}  ✕")
            chat_box.config(state=tk.NORMAL)
            chat_box.insert(tk.END, f"\n  📎 Archivo cargado: {nombre}\n", "file_tag")
            chat_box.config(state=tk.DISABLED)
            chat_box.see(tk.END)
    else:
        chat_box.config(state=tk.NORMAL)
        chat_box.insert(tk.END, f"\n  ⚠ No se pudo leer: {nombre}\n", "system_text")
        chat_box.config(state=tk.DISABLED)
        chat_box.see(tk.END)

def reproducir_cancion(nombre, silencio=False):
    def responder(texto):
        agregar_mensaje("JARVIS", texto)
        if not silencio: hablar(texto)
        historial_chat.append({"role": "assistant", "content": texto})
        guardar_historial()
    try:
        youtube = build("youtube", "v3", developerKey=YOUTUBE_API_KEY)
        resultados = youtube.search().list(q=nombre, part="snippet", maxResults=1, type="video").execute()
        items = resultados.get("items", [])
        if not items:
            reproducir_sfx("error")
            responder("No pude encontrar esa canción. ¿Podría repetir el nombre?")
            return False
        video_id = items[0]["id"]["videoId"]
        titulo = items[0]["snippet"]["title"]
        webbrowser.open(f"https://www.youtube.com/watch?v={video_id}")
        reproducir_sfx("success")
        responder(f"Reproduciendo {titulo} en YouTube, Pedro.")
        historial_chat.append({"role": "system", "content": f"[ACCIÓN] Se reprodujo en YouTube: '{titulo}' (búsqueda: '{nombre}')"})
        guardar_historial()
        return True
    except:
        reproducir_sfx("error")
        responder("No se ha encontrado la canción pedida.")
        return False


# ─── CREACIÓN DE DOCUMENTOS Y NUEVAS FUNCIONES ────────────────────────────────

def crear_gmail(instrucciones):
    """Extrae datos de la voz/texto y abre un borrador de Gmail pre-rellenado."""
    set_estado("PREPARANDO GMAIL...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere redactar un correo electrónico basándose en esta instrucción:
{instrucciones}

Analiza la instrucción y extrae o genera el destinatario, el asunto y el cuerpo del correo.
Responde SOLO en JSON con este formato exacto, sin texto extra ni markdown:
{{
  "destinatario": "correo_destino@gmail.com",
  "asunto": "Asunto del correo",
  "cuerpo": "Cuerpo completo y formal del correo."
}}"""
    try:
        raw = preguntar_ia(prompt, max_tokens=1024).strip()
        raw = raw[raw.find("{"):raw.rfind("}")+1]
        data = json.loads(raw)
        
        dest = data.get("destinatario", "")
        # Si el usuario no dictó un dominio, evitamos errores básicos
        if "@" not in dest and dest != "":
            dest = ""

        asunto_q = urllib.parse.quote(data.get("asunto", "Sin asunto"))
        cuerpo_q = urllib.parse.quote(data.get("cuerpo", ""))
        
        # Generar enlace mailto para Google Mail (Abre ventana de redacción)
        url_gmail = f"https://mail.google.com/mail/?view=cm&fs=1&to={dest}&su={asunto_q}&body={cuerpo_q}"
        webbrowser.open(url_gmail)
        reproducir_sfx("success")
        return "Borrador preparado en tu navegador, señor."
    except Exception as e:
        reproducir_sfx("error")
        return f"Error al generar el correo: {str(e)}"

def abrir_programa_o_ruta(consulta):
    """Usa PyAutoGUI para presionar Win + escribir el programa, simulando la búsqueda de Windows."""
    try:
        programa = consulta.lower().replace("abre el programa", "").replace("abrir programa", "").strip()
        pyautogui.press('win')
        time.sleep(0.5)
        pyautogui.write(programa)
        time.sleep(0.8)
        pyautogui.press('enter')
        reproducir_sfx("success")
        return f"Ejecutando {programa} mediante el sistema principal."
    except Exception as e:
        reproducir_sfx("error")
        return "Hubo un error al intentar abrir la aplicación en su PC."

def cambiar_salida_audio():
    """Abre el panel de sonido clásico de Windows (uso genérico, sin nombre de dispositivo)."""
    try:
        subprocess.run(["control", "mmsys.cpl,,0"])
        reproducir_sfx("success")
        return "Panel de dispositivos de reproducción abierto en pantalla."
    except Exception:
        reproducir_sfx("error")
        return "No pude acceder al panel de sonido del sistema."


# ─── CAMBIO REAL DE SALIDA DE AUDIO (con apodos guardados) ─────────────────
# Requiere el módulo de PowerShell "AudioDeviceCmdlets" (una sola vez):
#   Install-Module -Name AudioDeviceCmdlets -Force -Scope CurrentUser

def _cargar_apodos_audio():
    if os.path.exists(SALIDAS_AUDIO_PATH):
        try:
            with open(SALIDAS_AUDIO_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            _log_error("_cargar_apodos_audio()", e)
    return {}

def _guardar_apodos_audio(apodos):
    try:
        os.makedirs(os.path.dirname(SALIDAS_AUDIO_PATH), exist_ok=True)
        with open(SALIDAS_AUDIO_PATH, "w", encoding="utf-8") as f:
            json.dump(apodos, f, ensure_ascii=False, indent=2)
    except Exception as e:
        _log_error("_guardar_apodos_audio()", e)

def _listar_dispositivos_reproduccion():
    """Devuelve la lista de dispositivos de salida (Playback) reales del sistema."""
    try:
        cmd = [
            "powershell", "-NoProfile", "-Command",
            "Get-AudioDevice -List | Where-Object {$_.Type -eq 'Playback'} | "
            "Select-Object Index, ID, Name, Default | ConvertTo-Json -Compress"
        ]
        resultado = subprocess.run(cmd, capture_output=True, text=True, timeout=10, creationflags=NO_WINDOW)
        salida = resultado.stdout.strip()
        if not salida:
            return []
        datos = json.loads(salida)
        if isinstance(datos, dict):  # PowerShell no envuelve en lista si hay un solo resultado
            datos = [datos]
        return datos
    except Exception as e:
        _log_error("_listar_dispositivos_reproduccion()", e)
        return []

def _buscar_dispositivo_por_texto(nombre_pedido, dispositivos):
    objetivo = _normalizar_texto(nombre_pedido)
    for d in dispositivos:
        nombre_normalizado = _normalizar_texto(d.get("Name", ""))
        if objetivo and (objetivo in nombre_normalizado or nombre_normalizado in objetivo):
            return d
    return None

def _aplicar_dispositivo_audio(id_dispositivo):
    try:
        id_escapado = id_dispositivo.replace("'", "''")
        cmd = ["powershell", "-NoProfile", "-Command", f"Set-AudioDevice -ID '{id_escapado}'"]
        resultado = subprocess.run(cmd, capture_output=True, text=True, timeout=10, creationflags=NO_WINDOW)
        return resultado.returncode == 0
    except Exception as e:
        _log_error("_aplicar_dispositivo_audio()", e)
        return False

def _extraer_nombre_salida_pedida(voz_normalizada):
    """Extrae lo que el usuario pidió después de 'salida de audio/sonido a ...'"""
    m = re.search(r"salida de (?:audio|sonido)\s+a\s+(.+)", voz_normalizada)
    if m:
        return m.group(1).strip(" .,")
    return None

def flujo_cambiar_salida_audio(nombre_pedido):
    """Cambia la salida de audio a lo que el usuario pidió: primero revisa si ya
    existe un apodo guardado con ese nombre; si no, busca entre los dispositivos
    reales del sistema y, si logra cambiarlo, pregunta si desea guardarle un apodo."""
    nombre_normalizado = _normalizar_texto(nombre_pedido)
    apodos = _cargar_apodos_audio()

    # 1) ¿Ya existe un apodo guardado con ese nombre?
    if nombre_normalizado in apodos:
        disp = apodos[nombre_normalizado]
        if _aplicar_dispositivo_audio(disp["id"]):
            agregar_mensaje("JARVIS", f"Salida de audio cambiada a {disp['nombre_real']} (\"{nombre_pedido}\").")
            hablar(f"Listo, cambié la salida de sonido a {nombre_pedido}")
        else:
            agregar_mensaje("JARVIS", f"No pude cambiar a la salida guardada como \"{nombre_pedido}\". Puede que ese dispositivo ya no esté disponible.")
            hablar(f"No pude cambiar a {nombre_pedido}. Es posible que ya no esté conectado.")
        return

    # 2) Buscar entre los dispositivos reales del sistema
    dispositivos = _listar_dispositivos_reproduccion()
    if not dispositivos:
        agregar_mensaje("JARVIS", "No pude obtener la lista de dispositivos de audio. Verifica que el módulo AudioDeviceCmdlets de PowerShell esté instalado.")
        hablar("No pude acceder a los dispositivos de audio del sistema.")
        return

    encontrado = _buscar_dispositivo_por_texto(nombre_pedido, dispositivos)
    if not encontrado:
        nombres_disponibles = ", ".join(d.get("Name", "?") for d in dispositivos)
        agregar_mensaje("JARVIS", f"No encontré ningún dispositivo de audio llamado \"{nombre_pedido}\". Disponibles: {nombres_disponibles}")
        hablar(f"No encontré una salida de audio llamada {nombre_pedido}")
        return

    if not _aplicar_dispositivo_audio(encontrado.get("ID", "")):
        agregar_mensaje("JARVIS", f"Encontré \"{encontrado.get('Name')}\" pero no pude activarlo como salida predeterminada.")
        hablar("Encontré el dispositivo pero no pude activarlo.")
        return

    nombre_real = encontrado.get("Name", nombre_pedido)
    agregar_mensaje("JARVIS", f"Salida de audio cambiada a {nombre_real}.")
    hablar(f"Listo, cambié la salida de sonido a {nombre_real}. ¿Quisiera ponerle un nombre a esta salida?")

    set_estado("ESCUCHANDO...")
    reproducir_sfx("listening")
    respuesta = _normalizar_texto(escuchar())

    PALABRAS_SI = ["si", "claro", "dale", "va", "ponle", "nombre"]
    if respuesta and _contiene_alguna(respuesta, PALABRAS_SI):
        m = re.search(r"nombre de\s+(.+)", respuesta)
        apodo = m.group(1).strip(" .,") if m else ""

        if not apodo:
            hablar("¿Qué nombre quiere ponerle?")
            set_estado("ESCUCHANDO...")
            reproducir_sfx("listening")
            apodo = escuchar().strip()

        if apodo:
            apodo_normalizado = _normalizar_texto(apodo)
            apodos[apodo_normalizado] = {"id": encontrado.get("ID"), "nombre_real": nombre_real}
            _guardar_apodos_audio(apodos)
            hablar(f"Perfecto, la recordaré como {apodo}")
            agregar_mensaje("JARVIS", f"Guardé esta salida de audio con el nombre \"{apodo}\".")
        else:
            hablar("No logré escuchar el nombre, no guardé el apodo.")
    else:
        hablar("Entendido, no le pondré nombre.")

def reporte_clima():
    """Obtiene el clima actual de Chiclayo."""
    try:
        res = requests.get("https://wttr.in/Chiclayo?format=%C+%t")
        clima_actual = res.text.strip()
        # Traducir pequeñas respuestas comunes si es necesario, o devolver directo
        return f"Las condiciones actuales en Chiclayo son: {clima_actual}."
    except Exception:
        return "Los servidores meteorológicos no están respondiendo en este momento."

def modo_gaming():
    """Cierra procesos pesados para liberar RAM y ajusta volumen."""
    subprocess.run(["taskkill", "/f", "/im", "chrome.exe"], creationflags=NO_WINDOW)
    subprocess.run(["taskkill", "/f", "/im", "opera.exe"], creationflags=NO_WINDOW)
    ajustar_volumen("set", 100)
    reproducir_sfx("success")
    return "Modo Gaming inicializado. Aplicaciones de fondo purgadas, memoria RAM liberada y volumen al máximo. Listo para jugar."

def crear_word(instrucciones):
    set_estado("GENERANDO WORD...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere crear un documento Word con estas instrucciones:
{instrucciones}
Responde SOLO en JSON con este formato exacto, sin texto extra ni markdown:
{{
  "titulo": "Título del documento",
  "subtitulo": "Subtítulo opcional o vacío",
  "parrafos": [
    {{ "tipo": "titulo_seccion", "texto": "Nombre de la sección" }},
    {{ "tipo": "parrafo", "texto": "Texto normal del párrafo." }},
    {{ "tipo": "lista", "items": ["item 1", "item 2", "item 3"] }},
    {{ "tipo": "tabla", "encabezados": ["Col1", "Col2"], "filas": [["a", "b"], ["c", "d"]] }},
    {{ "tipo": "salto_linea" }}
  ]
}}"""

    import json, re, os
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    try:
        raw = preguntar_ia(prompt, max_tokens=4096).strip()
        start = raw.find("{")
        end = raw.rfind("}") + 1
        raw = raw[start:end]
        raw = re.sub(r',\s*}', '}', raw)
        raw = re.sub(r',\s*]', ']', raw)
        data = json.loads(raw)
    except Exception as e:
        reproducir_sfx("error")
        raise e

    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    titulo_p = doc.add_paragraph()
    titulo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = titulo_p.add_run(data["titulo"])
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3E, 0x7A)
    run.font.name = "Calibri"

    if data.get("subtitulo"):
        sub_p = doc.add_paragraph()
        sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_p.add_run(data["subtitulo"])
        sub_run.font.size = Pt(13)
        sub_run.font.italic = True
        sub_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    doc.add_paragraph()

    for bloque in data.get("parrafos", []):
        tipo = bloque.get("tipo", "parrafo")
        if tipo == "titulo_seccion":
            h = doc.add_paragraph()
            h_run = h.add_run(bloque["texto"])
            h_run.font.size = Pt(14)
            h_run.font.bold = True
            h_run.font.color.rgb = RGBColor(0x1F, 0x3E, 0x7A)
            h_run.font.name = "Calibri"
        elif tipo == "parrafo":
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.2)
            run = p.add_run(bloque["texto"])
            run.font.size = Pt(11)
            run.font.name = "Calibri"
        elif tipo == "lista":
            for item in bloque.get("items", []):
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(item)
                run.font.size = Pt(11)
                run.font.name = "Calibri"
        elif tipo == "tabla":
            encabezados = bloque.get("encabezados", [])
            filas = bloque.get("filas", [])
            if encabezados:
                tabla = doc.add_table(rows=1 + len(filas), cols=len(encabezados))
                tabla.style = "Table Grid"
                for i, enc in enumerate(encabezados):
                    cell = tabla.rows[0].cells[i]
                    cell.text = enc
                    if cell.paragraphs[0].runs:
                        cell.paragraphs[0].runs[0].font.bold = True
                        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    shading = OxmlElement("w:shd")
                    shading.set(qn("w:fill"), "1F3E7A")
                    shading.set(qn("w:color"), "auto")
                    shading.set(qn("w:val"), "clear")
                    cell._tc.get_or_add_tcPr()
                    cell._tc.tcPr.append(shading)
                for r, fila in enumerate(filas):
                    for c, dato in enumerate(fila):
                        tabla.rows[r+1].cells[c].text = str(dato)
                doc.add_paragraph()

    os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
    nombre_limpio = re.sub(r'[\\/:*?"<>|]', '', data["titulo"]).replace(' ', '_')
    ruta = f"C:/Jarvis/Documentos/{nombre_limpio}.docx"
    doc.save(ruta)
    os.startfile(ruta)
    reproducir_sfx("success")
    return data["titulo"]

def crear_pptx(instrucciones):
    set_estado("GENERANDO POWERPOINT...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere crear una presentación PowerPoint con estas instrucciones:
{instrucciones}
Responde SOLO en JSON con este formato exacto, sin texto extra:
{{
  "titulo": "Título de la presentación",
  "subtitulo": "Subtítulo breve",
  "tema": "oscuro",
  "color_primario": "6C8CFF",
  "color_secundario": "A78BFA",
  "color_fondo": "0D0D1A",
  "diapositivas": [
    {{
      "titulo": "Título diapositiva",
      "layout": "puntos",
      "contenido": ["punto 1", "punto 2"]
    }}
  ]
}}"""

    import json, re, os
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    try:
        raw = preguntar_ia(prompt, max_tokens=2048).strip()
        raw = raw[raw.find("{"):raw.rfind("}")+1]
        data = json.loads(raw)
    except Exception as e:
        reproducir_sfx("error")
        raise e

    def hex_to_rgb(h):
        h = h.lstrip("#").strip()
        if len(h) != 6: h = "6C8CFF"
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    C_BG    = hex_to_rgb(data.get("color_fondo",    "0D0D1A"))
    C_PRI   = hex_to_rgb(data.get("color_primario",  "6C8CFF"))
    C_SEC   = hex_to_rgb(data.get("color_secundario","A78BFA"))
    C_WHITE = RGBColor(0xFF,0xFF,0xFF)
    C_DIM   = RGBColor(0xB0,0xB0,0xCC)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    def blank(): return prs.slides.add_slide(prs.slide_layouts[6])

    def rect(slide, x, y, w, h, color):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background(); return s

    def txt(slide, text, x, y, w, h, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
        r.font.name = "Segoe UI"

    slide = blank()
    rect(slide, 0, 0, W, H, C_BG)
    txt(slide, data["titulo"], Inches(0.55), Inches(2.0), Inches(9), Inches(1.4), 46, C_WHITE, bold=True)
    txt(slide, data.get("subtitulo","Generado por J.A.R.V.I.S"), Inches(0.55), Inches(3.6), Inches(8), Inches(0.7), 18, C_DIM, italic=True)

    os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
    nombre_limpio = re.sub(r'[\\/:*?"<>|]', '', data["titulo"]).replace(' ', '_')
    ruta = f"C:/Jarvis/Documentos/{nombre_limpio}.pptx"
    prs.save(ruta)
    os.startfile(ruta)
    reproducir_sfx("success")
    return data["titulo"]

def crear_excel(instrucciones):
    set_estado("GENERANDO EXCEL...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere crear una hoja de Excel con estas instrucciones:
{instrucciones}
Responde SOLO en JSON con este formato exacto, sin texto extra:
{{
  "titulo": "Nombre del archivo",
  "hojas": [
    {{
      "nombre": "Hoja1",
      "encabezados": ["Col1", "Col2"],
      "filas": [["dato1", "dato2"]]
    }}
  ]
}}"""

    import json, re, os
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    try:
        raw = preguntar_ia(prompt, max_tokens=2048).strip()
        raw = raw[raw.find("{"):raw.rfind("}")+1]
        data = json.loads(raw)
    except Exception as e:
        reproducir_sfx("error")
        raise e

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    for hoja_data in data["hojas"]:
        ws = wb.create_sheet(title=hoja_data["nombre"])
        encabezados = hoja_data["encabezados"]
        filas = hoja_data["filas"]

        for col, enc in enumerate(encabezados, 1):
            cell = ws.cell(row=1, column=col, value=enc)
            cell.font = Font(bold=True, color="FFFFFF", name="Calibri", size=12)
            cell.fill = PatternFill("solid", fgColor="1F3E7A")

        for r, fila in enumerate(filas, 2):
            for c, dato in enumerate(fila, 1):
                ws.cell(row=r, column=c, value=dato)

    os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
    nombre_limpio = re.sub(r'[\\/:*?"<>|]', '', data["titulo"]).replace(' ', '_')
    ruta = f"C:/Jarvis/Documentos/{nombre_limpio}.xlsx"
    wb.save(ruta)
    os.startfile(ruta)
    reproducir_sfx("success")
    return data["titulo"]


def crear_archivo_py(instrucciones):
    set_estado("GENERANDO PYTHON...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere crear un archivo Python con estas instrucciones:
{instrucciones}
Responde SOLO en JSON con este formato exacto, sin texto extra:
{{
  "nombre": "nombre_del_archivo",
  "descripcion": "Breve descripción",
  "codigo": "código python"
}}"""
    import json, re, os
    try:
        raw = preguntar_ia(prompt, max_tokens=2048).strip()
        raw = raw[raw.find("{"):raw.rfind("}")+1]
        data = json.loads(raw)
    except Exception as e:
        reproducir_sfx("error")
        raise e

    os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
    nombre_limpio = re.sub(r'[\\/:*?"<>|]', '', data["nombre"]).replace(' ', '_')
    ruta = f"C:/Jarvis/Documentos/{nombre_limpio}.py"
    with open(ruta, "w", encoding="utf-8") as f:
        f.write(data["codigo"])
    os.startfile(ruta)
    reproducir_sfx("success")
    return data["nombre"], data["descripcion"]

def pb():
    if not PB_API_KEY:
        raise RuntimeError("PUSHBULLET_API_KEY no configurada.")
    return Pushbullet(PB_API_KEY)

def enviar_notificacion(titulo, mensaje):
    try:
        pb().push_note(titulo, mensaje)
        reproducir_sfx("success")
        return True
    except Exception as e:
        reproducir_sfx("error")
        return False

def recordatorio(mensaje, hora_str):
    import time
    try:
        hoy = datetime.now()
        h, m = map(int, hora_str.replace(".", ":").split(":"))
        objetivo = hoy.replace(hour=h, minute=m, second=0, microsecond=0)
        if objetivo < hoy:
            objetivo = objetivo + timedelta(days=1)
        segundos = (objetivo - datetime.now()).total_seconds()

        def _esperar_y_enviar():
            time.sleep(segundos)
            enviar_notificacion("⏰ Recordatorio de J.A.R.V.I.S", mensaje)

        threading.Thread(target=_esperar_y_enviar, daemon=True).start()
        return objetivo.strftime("%H:%M")
    except Exception:
        reproducir_sfx("error")
        return None

def abrir_plataforma_con_login():
    try:
        options = Options()
        options.add_argument("--start-maximized")
        options.add_argument(r"--user-data-dir=C:\Users\MSi\AppData\Local\Google\Chrome\User Data")
        options.add_argument("--profile-directory=Profile 3")
        options.add_experimental_option("excludeSwitches", ["enable-automation"])
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=options)
        driver.get("https://plataforma.colegiocima.edu.pe/CampusVirtual/#/login")
    except Exception:
        webbrowser.open("https://plataforma.colegiocima.edu.pe/CampusVirtual/#/login")


def ajustar_volumen(accion, valor=None):
    try:
        from pycaw.pycaw import AudioUtilities
        devices = AudioUtilities.GetSpeakers()
        volume = devices.EndpointVolume
        vol_actual = round(volume.GetMasterVolumeLevelScalar() * 100)
        if accion == "set" and valor is not None:
            nuevo = max(0, min(100, valor))
        elif accion == "subir":
            nuevo = min(100, vol_actual + 5)
        elif accion == "bajar":
            nuevo = max(0, vol_actual - 5)
        else:
            return vol_actual
        volume.SetMasterVolumeLevelScalar(nuevo / 100, None)
        return nuevo
    except Exception:
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# ─── FUNCIONES IA AVANZADAS v2.0 ────────────────────────────────────────────
# ═══════════════════════════════════════════════════════════════════════════════

# ─── ESCANEO DE DOCUMENTOS CON IA ──────────────────────────────────────────
def escanear_documento_ia(ruta_archivo=None):
    """Analiza un documento cargado usando IA y devuelve un resumen inteligente."""
    contenido = None
    if ruta_archivo:
        contenido = leer_archivo(ruta_archivo)
    elif archivo_cargado.get("contenido"):
        contenido = archivo_cargado["contenido"]
    else:
        return "No hay ningún documento cargado. Carga uno primero con el botón +."
    if not contenido:
        return "No pude leer el contenido del archivo."
    set_estado("ESCANEANDO DOCUMENTO...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. Analiza este documento y proporciona:
1. Resumen ejecutivo (3-5 oraciones)
2. Puntos clave principales
3. Palabras o conceptos importantes
4. Posibles errores o mejoras sugeridas (si aplica)

DOCUMENTO:
{contenido[:5000]}"""
    try:
        respuesta = preguntar_ia(prompt, max_tokens=2048)
        reproducir_sfx("success")
        return f"📋 Análisis del documento:\n\n{respuesta}"
    except Exception as e:
        reproducir_sfx("error")
        return f"No pude analizar el documento: {e}"


# ─── CREAR ARCHIVOS CON IA ─────────────────────────────────────────────────
def crear_archivo_ia(instrucciones):
    """Crea cualquier tipo de archivo basándose en instrucciones del usuario."""
    set_estado("CREANDO ARCHIVO CON IA...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. El usuario quiere crear un archivo. Analiza esta instrucción:
{instrucciones}

Determina el tipo de archivo más apropiado y genera el contenido completo.
Responde SOLO en JSON:
{{
  "nombre": "nombre_del_archivo",
  "tipo": "py|txt|html|css|js|json|md|csv",
  "contenido": "contenido completo del archivo",
  "descripcion": "breve descripción de lo creado"
}}"""
    try:
        raw = preguntar_ia(prompt, max_tokens=4096).strip()
        raw = raw[raw.find("{"):raw.rfind("}")+1]
        data = json.loads(raw)
        os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
        nombre_limpio = re.sub(r'[\\/:*?"<>|]', '', data["nombre"]).replace(' ', '_')
        ext = data.get("tipo", "txt")
        ruta = f"C:/Jarvis/Documentos/{nombre_limpio}.{ext}"
        with open(ruta, "w", encoding="utf-8") as f:
            f.write(data["contenido"])
        reproducir_sfx("success")
        return f"✅ Archivo creado: {nombre_limpio}.{ext}\n📝 {data.get('descripcion', 'Archivo generado por IA')}\n📁 Ubicación: {ruta}"
    except Exception as e:
        reproducir_sfx("error")
        return f"No pude crear el archivo: {e}"


# ─── MODIFICAR ARCHIVOS CON IA ─────────────────────────────────────────────
def modificar_archivo_ia(ruta_archivo, instrucciones):
    """Lee un archivo existente, lo modifica con IA y lo guarda."""
    contenido = leer_archivo(ruta_archivo)
    if not contenido:
        return f"No pude leer el archivo: {ruta_archivo}"
    set_estado("MODIFICANDO ARCHIVO CON IA...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. Modifica este archivo según las instrucciones.
Devuelve SOLO el contenido modificado completo, sin explicaciones extra.

ARCHIVO ACTUAL:
{contenido[:4000]}

INSTRUCCIONES DE MODIFICACIÓN:
{instrucciones}"""
    try:
        nuevo_contenido = preguntar_ia(prompt, max_tokens=4096).strip()
        # Crear backup
        backup_dir = "C:/Jarvis/backups"
        os.makedirs(backup_dir, exist_ok=True)
        nombre = os.path.basename(ruta_archivo)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_path = f"{backup_dir}/{timestamp}_{nombre}"
        import shutil
        shutil.copy2(ruta_archivo, backup_path)
        with open(ruta_archivo, "w", encoding="utf-8") as f:
            f.write(nuevo_contenido)
        reproducir_sfx("success")
        return f"✅ Archivo modificado: {nombre}\n💾 Backup guardado en: {backup_path}"
    except Exception as e:
        reproducir_sfx("error")
        return f"No pude modificar el archivo: {e}"


# ─── LEER ARCHIVOS POR VOZ ─────────────────────────────────────────────────
def leer_archivo_por_voz(nombre_archivo):
    """Busca un archivo por nombre y lee su contenido con IA."""
    nombre_lower = nombre_archivo.lower().strip()
    rutas_busqueda = [
        f"C:/Jarvis/{nombre_archivo}",
        f"C:/Jarvis/Documentos/{nombre_archivo}",
        f"C:/Jarvis/config/{nombre_archivo}",
    ]
    # Buscar recursivamente en C:/Jarvis
    for root, dirs, files in os.walk("C:/Jarvis"):
        for f in files:
            if nombre_lower in f.lower():
                rutas_busqueda.insert(0, os.path.join(root, f))
    for ruta in rutas_busqueda:
        if os.path.exists(ruta):
            contenido = leer_archivo(ruta)
            if contenido:
                return f"📄 Contenido de {os.path.basename(ruta)}:\n\n{contenido[:3000]}"
    return f"No encontré ningún archivo llamado '{nombre_archivo}'."


# ─── ANALIZAR CÓDIGO CON IA ─────────────────────────────────────────────────
def analizar_codigo_ia(ruta_archivo=None):
    """Analiza código fuente con IA: bugs, mejoras, explicación."""
    contenido = None
    if ruta_archivo:
        contenido = leer_archivo(ruta_archivo)
    elif archivo_cargado.get("contenido"):
        contenido = archivo_cargado["contenido"]
    else:
        return "No hay código cargado. Carga un archivo primero."
    set_estado("ANALIZANDO CÓDIGO...")
    reproducir_sfx("processing")
    prompt = f"""Eres J.A.R.V.I.S. Analiza este código y proporciona:
1. Resumen de qué hace
2. Posibles bugs o errores
3. Mejoras de rendimiento o diseño
4. Complejidad del código

CÓDIGO:
{contenido[:5000]}"""
    try:
        respuesta = preguntar_ia(prompt, max_tokens=2048)
        reproducir_sfx("success")
        return f"🔍 Análisis de código:\n\n{respuesta}"
    except Exception as e:
        return f"No pude analizar el código: {e}"


# ─── EJECUTAR CÓDIGO PYTHON DE FORMA SEGURA ────────────────────────────────
def ejecutar_codigo_python(codigo):
    """Ejecuta código Python en un entorno temporal y captura la salida."""
    set_estado("EJECUTANDO CÓDIGO...")
    reproducir_sfx("processing")
    try:
        resultado = []
        def _ejecutar():
            import contextlib
            old_stdout = sys.stdout
            old_stderr = sys.stderr
            redirected_output = io.StringIO()
            redirected_error = io.StringIO()
            sys.stdout = redirected_output
            sys.stderr = redirected_error
            try:
                exec(codigo, {"__builtins__": __builtins__})
            except Exception as e:
                print(f"Error: {e}", file=sys.stderr)
            finally:
                sys.stdout = old_stdout
                sys.stderr = old_stderr
                resultado.append(redirected_output.getvalue())
                resultado.append(redirected_error.getvalue())
        _ejecutar()
        salida = resultado[0] if resultado and resultado[0] else "(sin salida)"
        errores = resultado[1] if len(resultado) > 1 and resultado[1] else ""
        reproducir_sfx("success")
        respuesta = f"📤 Salida:\n{salida}"
        if errores:
            respuesta += f"\n⚠️ Errores:\n{errores}"
        return respuesta
    except Exception as e:
        reproducir_sfx("error")
        return f"Error al ejecutar: {e}"


# ─── BUSCAR ARCHIVOS EN EL SISTEMA ─────────────────────────────────────────
def buscar_archivos(consulta):
    """Busca archivos en el sistema por nombre."""
    resultados = []
    for root, dirs, files in os.walk("C:/Jarvis"):
        # Ignorar carpetas de versiones y .git
        dirs[:] = [d for d in dirs if d not in ['Versiones', '.git', '__pycache__', 'logs']]
        for f in files:
            if consulta.lower() in f.lower():
                resultados.append(os.path.join(root, f))
    if not resultados:
        return f"No encontré archivos con '{consulta}' en C:/Jarvis."
    respuesta = f"🔍 Archivos encontrados ({len(resultados)}):\n"
    for ruta in resultados[:10]:
        tamano = os.path.getsize(ruta) / 1024
        respuesta += f"  📄 {os.path.basename(ruta)} ({tamano:.1f} KB) — {os.path.dirname(ruta)}\n"
    return respuesta


# ─── LISTAR CARPETA ────────────────────────────────────────────────────────
def listar_carpeta(ruta="C:/Jarvis"):
    """Lista el contenido de una carpeta."""
    if not os.path.exists(ruta):
        return f"La carpeta {ruta} no existe."
    items = os.listdir(ruta)
    carpetas = [i for i in items if os.path.isdir(os.path.join(ruta, i))]
    archivos = [i for i in items if os.path.isfile(os.path.join(ruta, i))]
    respuesta = f"📁 {ruta}:\n"
    if carpetas:
        respuesta += f"  📂 Carpetas ({len(carpetas)}): {', '.join(carpetas[:15])}\n"
    if archivos:
        respuesta += f"  📄 Archivos ({len(archivos)}): {', '.join(archivos[:15])}\n"
    return respuesta


def mostrar_toast(mensaje, tipo="info", duracion_ms=3500):
    """Muestra una notificación flotante tipo toast en la esquina superior derecha."""
    colores = {"info": ACCENT, "success": GREEN, "error": RED, "warning": AMBER}
    iconos = {"info": "ℹ", "success": "✓", "error": "✗", "warning": "⚠"}
    color = colores.get(tipo, ACCENT)
    icono = iconos.get(tipo, "ℹ")

    def _crear():
        global toast_widgets
        toast = tk.Frame(ventana, bg=SURFACE3, highlightthickness=1, highlightbackground=color)
        toast.place(relx=1.0, rely=0.0, anchor="ne", x=-20, y=20 + len(toast_widgets) * 50)
        tk.Label(toast, text=icono, font=("Segoe UI", 12, "bold"), fg=color, bg=SURFACE3).pack(side=tk.LEFT, padx=(12, 6), pady=8)
        tk.Label(toast, text=mensaje, font=("Segoe UI", 10), fg=TEXT, bg=SURFACE3, wraplength=280, justify="left").pack(side=tk.LEFT, padx=(0, 12), pady=8)
        toast_widgets.append(toast)

        def _cerrar():
            try:
                toast.destroy()
                if toast in toast_widgets: toast_widgets.remove(toast)
            except: pass

        ventana.after(duracion_ms, _cerrar)
        # Animar entrada
        try: toast.place_configure(x=-300)
        except: pass
        _animar_toast_entrada(toast, -300, -20, 0)

    def _animar_toast_entrada(t, x_inicio, x_fin, paso):
        total_pasos = 12
        if paso > total_pasos:
            try: t.place_configure(x=x_fin)
            except: pass
            return
        t_obj = paso / total_pasos
        x = x_inicio + (x_fin - x_inicio) * (1 - math.pow(1 - t_obj, 3))
        try: t.place_configure(x=int(x))
        except: pass
        ventana.after(18, lambda: _animar_toast_entrada(t, x_inicio, x_fin, paso + 1))

    ui(_crear)


# ─── CALCULADORA AVANZADA (ELIMINADA — ver funciones IA arriba) ───────────
# Estas funciones fueron reemplazadas por las funciones IA.
def _eliminar_bloque_antiguo():
    pass

# Las funciones de calculadora, contraseñas, screenshots, clipboard,
# notas, timer, pomodoro, info sistema, traducción y resumen fueron
# reemplazadas por las funciones IA superiores.
# BLKDEL_INICIO
    expr = expr.replace("x", "*").replace("por", "*").replace("dividido", "/")
    expr = expr.replace("elevado a", "**").replace("potencia", "**")
    expr = expr.replace("raiz cuadrada", "math.sqrt").replace("raíz cuadrada", "math.sqrt")
    expr = expr.replace("seno", "math.sin").replace("coseno", "math.cos")
    expr = expr.replace("tangente", "math.tan")
    expr = expr.replace("logaritmo", "math.log10").replace("logaritmo natural", "math.log")
    expr = expr.replace("factorial", "math.factorial")
    expr = expr.replace("valor absoluto", "abs")
    expr = expr.replace("pi", "math.pi")
    expr = re.sub(r'[^0-9+\-*/().,e \t\n]', '', expr)
    expr = expr.strip()
    if not expr:
        return "No se detectó una expresión matemática válida."
    try:
        resultado = eval(expr, {"__builtins__": {}, "math": math})
        if isinstance(resultado, float):
            if resultado == int(resultado) and abs(resultado) < 1e15:
                resultado = int(resultado)
            else:
                resultado = round(resultado, 10)
        return f"El resultado es: {resultado}"
    except Exception:
        # Si falla directo, pedir a la IA que resuelva
        try:
            respuesta = preguntar_ia(f"Resuelve esta operación matemática y responde SOLO con el número resultado: {expr}", max_tokens=64)
            return f"El resultado es: {respuesta.strip()}"
        except:
            return "No pude resolver esa expresión matemática."


# ─── GENERADOR DE CONTRASEÑAS SEGURAS ──────────────────────────────────────
def generar_contrasena(longitud=16, usar_mayusculas=True, usar_numeros=True, usar_simbolos=True):
    """Genera una contraseña criptográficamente segura."""
    caracteres = string.ascii_lowercase
    if usar_mayusculas: caracteres += string.ascii_uppercase
    if usar_numeros: caracteres += string.digits
    if usar_simbolos: caracteres += "!@#$%^&*()-_=+[]{}|;:,.<>?"
    longitud = max(8, min(128, longitud))
    contrasena = ''.join(secrets.choice(caracteres) for _ in range(longitud))
    # Calcular entropía
    entropia = len(caracteres) ** longitud
    bits_entropia = math.log2(entropia)
    nivel = "Excelente" if bits_entropia >= 80 else "Buena" if bits_entropia >= 60 else "Moderada"
    return f"🔒 Contraseña generada: {contrasena}\n   Longitud: {longitud} caracteres | Entropía: {bits_entropia:.0f} bits | Seguridad: {nivel}"


# ─── CAPTURA DE PANTALLA ───────────────────────────────────────────────────
def capturar_pantalla(region=None):
    """Captura la pantalla actual y la guarda como imagen."""
    try:
        os.makedirs("C:/Jarvis/Documentos", exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ruta = f"C:/Jarvis/Documentos/captura_{timestamp}.png"
        screenshot = pyautogui.screenshot(region=region)
        screenshot.save(ruta)
        mostrar_toast(f"Captura guardada: {timestamp}.png", "success")
        return ruta
    except Exception as e:
        _log_error("capturar_pantalla()", e)
        return None


# ─── CLIPBOARD MANAGER ──────────────────────────────────────────────────────
def copiar_al_portapapeles(texto):
    """Copia texto al portapapeles y lo registra en el historial."""
    try:
        ventana.clipboard_clear()
        ventana.clipboard_append(texto)
        # Registrar en historial
        entrada = {"texto": texto[:200], "timestamp": datetime.now().strftime("%H:%M:%S")}
        clipboard_history.insert(0, entrada)
        if len(clipboard_history) > clipboard_max:
            clipboard_history.pop()
        return True
    except: return False


def leer_portapapeles():
    """Lee el contenido actual del portapapeles."""
    try: return ventana.clipboard_get()
    except: return ""


def mostrar_historial_clipboard():
    if not clipboard_history:
        return "El historial del portapapeles está vacío."
    resultado = "📋 Historial del portapapeles:\n"
    for i, entrada in enumerate(clipboard_history[:10]):
        texto = entrada["texto"][:80]
        resultado += f"  {i+1}. [{entrada['timestamp']}] {texto}\n"
    return resultado


# ─── NOTAS RÁPIDAS ──────────────────────────────────────────────────────────
def cargar_notas():
    if os.path.exists(NOTAS_PATH):
        try:
            with open(NOTAS_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except: pass
    return []


def guardar_nota(contenido, titulo=""):
    """Guarda una nota rápida con timestamp."""
    notas = cargar_notas()
    nota = {
        "id": len(notas) + 1,
        "titulo": titulo or f"Nota {len(notas) + 1}",
        "contenido": contenido,
        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "favorita": False
    }
    notas.append(nota)
    try:
        os.makedirs(os.path.dirname(NOTAS_PATH), exist_ok=True)
        with open(NOTAS_PATH, "w", encoding="utf-8") as f:
            json.dump(notas, f, ensure_ascii=False, indent=2)
        mostrar_toast(f"Nota #{nota['id']} guardada", "success")
        return f"Nota #{nota['id']} guardada: {nota['titulo']}"
    except Exception as e:
        _log_error("guardar_nota()", e)
        return "Error al guardar la nota."


def buscar_notas(consulta):
    """Busca notas por contenido o título."""
    notas = cargar_notas()
    if not notas:
        return "No tienes notas guardadas aún."
    consulta_lower = consulta.lower()
    resultados = []
    for nota in notas:
        if consulta_lower in nota.get("titulo", "").lower() or consulta_lower in nota.get("contenido", "").lower():
            resultados.append(nota)
    if not resultados:
        resultados = notas[-5:]
    resultado = f"📝 Notas encontradas ({len(resultados)}):\n"
    for nota in resultados[-5:]:
        fav = "⭐" if nota.get("favorita") else "  "
        resultado += f"  {fav} #{nota['id']}: {nota['titulo']} ({nota['fecha']})\n     {nota['contenido'][:100]}...\n"
    return resultado


def listar_notas():
    """Lista todas las notas."""
    notas = cargar_notas()
    if not notas:
        return "📝 No tienes notas guardadas. Puedes crear una con 'guarda nota: ...'"
    resultado = f"📝 Tus notas ({len(notas)} total):\n"
    for nota in notas[-10:]:
        fav = "⭐" if nota.get("favorita") else "  "
        resultado += f"  {fav} #{nota['id']}: {nota['titulo']} — {nota['fecha']}\n"
    return resultado


def eliminar_nota(id_nota):
    """Elimina una nota por su ID."""
    notas = cargar_notas()
    notas_original = len(notas)
    notas = [n for n in notas if n.get("id") != id_nota]
    if len(notas) < notas_original:
        with open(NOTAS_PATH, "w", encoding="utf-8") as f:
            json.dump(notas, f, ensure_ascii=False, indent=2)
        return f"Nota #{id_nota} eliminada."
    return f"No encontré la nota #{id_nota}."


# ─── TEMPORIZADOR / CRONÓMETRO ──────────────────────────────────────────────
def iniciar_timer(segundos, mensaje="Tu tiempo ha terminado"):
    """Inicia un temporizador de cuenta regresiva."""
    global timer_corriendo, timer_hilo
    if timer_corriendo:
        return "Ya hay un temporizador activo. Usa 'detener temporizador' para cancelarlo."
    timer_corriendo = True
    def _esperar():
        global timer_corriendo
        time.sleep(segundos)
        timer_corriendo = False
        enviar_notificacion("⏰ Timer JARVIS", mensaje)
        hablar(f"Se acabó el tiempo. {mensaje}")
    timer_hilo = threading.Thread(target=_esperar, daemon=True)
    timer_hilo.start()
    mins = segundos // 60
    segs = segundos % 60
    tiempo_str = f"{mins} min {segs} seg" if mins > 0 else f"{segs} segundos"
    return f"⏱ Temporizador iniciado: {tiempo_str}. Te avisaré cuando termine."


def detener_timer():
    global timer_corriendo
    timer_corriendo = False
    return "Temporizador detenido."


# ─── POMODORO ───────────────────────────────────────────────────────────────
def iniciar_pomodoro(minutos_trabajo=25, minutos_descanso=5):
    """Inicia una sesión Pomodoro."""
    global pomodoro_activo, pomodoro_segundos, pomodoro_descanso
    pomodoro_activo = True
    pomodoro_segundos = minutos_trabajo * 60
    pomodoro_descanso = minutos_descanso * 60
    return f"🍅 Pomodoro iniciado: {minutos_trabajo} min trabajo → {minutos_descanso} min descanso. ¡Enfócate!"


def iniciar_sesion_pomodoro():
    """Inicia el temporizador Pomodoro en un hilo."""
    global pomodoro_activo
    if not pomodoro_activo:
        return
    def _ciclo_pomodoro():
        global pomodoro_activo
        while pomodoro_activo:
            # Fase de trabajo
            tiempo = pomodoro_segundos
            mins = tiempo // 60
            hablar(f"Pomodoro iniciado. Tienes {mins} minutos de trabajo. ¡Enfócate!")
            time.sleep(tiempo)
            if not pomodoro_activo: break
            enviar_notificacion("🍅 Pomodoro", "¡Se acabó el tiempo de trabajo! Toma un descanso.")
            hablar("Se acabó el tiempo de trabajo. ¡Toma un descanso de 5 minutos!")
            # Fase de descanso
            time.sleep(pomodoro_descanso)
            if not pomodoro_activo: break
            hablar("El descanso terminó. ¿Continuamos con otro Pomodoro? Di 'continuar pomodoro' para seguir.")
    threading.Thread(target=_ciclo_pomodoro, daemon=True).start()


def detener_pomodoro():
    global pomodoro_activo
    pomodoro_activo = False
    return "🍅 Pomodoro detenido. ¡Buen trabajo!"


# ─── INFORMACIÓN DEL SISTEMA ───────────────────────────────────────────────
def info_sistema_detallada():
    """Devuelve información detallada del sistema."""
    try:
        import platform
        info_parts = []
        info_parts.append(f"🖥️ Sistema: {platform.system()} {platform.release()}")
        info_parts.append(f"💻 PC: {platform.node()}")
        info_parts.append(f"🔧 Procesador: {platform.processor()}")
        if _psutil:
            ram = _psutil.virtual_memory()
            ram_uso = ram.used / (1024**3)
            ram_total = ram.total / (1024**3)
            ram_pct = ram.percent
            info_parts.append(f"📊 RAM: {ram_uso:.1f} / {ram_total:.1f} GB ({ram_pct}%)")
            disk = _psutil.disk_usage('C:/')
            disco_uso = disk.used / (1024**3)
            disco_total = disk.total / (1024**3)
            info_parts.append(f"💾 Disco C:: {disco_uso:.1f} / {disco_total:.1f} GB")
            if hasattr(_psutil, 'net_io_counters'):
                net = _psutil.net_io_counters()
                info_parts.append(f"🌐 Red: ↑{net.bytes_sent/(1024**2):.1f} MB  ↓{net.bytes_recv/(1024**2):.1f} MB")
            # Top 5 procesos
            procs = sorted(_psutil.process_iter(['pid', 'name', 'memory_percent']), key=lambda p: p.info.get('memory_percent', 0), reverse=True)
            info_parts.append("\n🏆 Top 5 procesos por RAM:")
            for p in procs[:5]:
                info_parts.append(f"   {p.info['name']}: {p.info.get('memory_percent', 0):.1f}%")
        else:
            info_parts.append("\n⚠️ psutil no instalado — instálalo con: pip install psutil")
        return "\n".join(info_parts)
    except Exception as e:
        return f"Error al obtener info del sistema: {e}"


# ─── TRADUCCIÓN RÁPIDA ──────────────────────────────────────────────────────
def traducir_texto(texto, idioma_destino="inglés"):
    """Traduce texto usando IA."""
    try:
        respuesta = preguntar_ia(
            f"Traduce el siguiente texto al {idioma_destino}. Responde SOLO con la traducción, nada más:\n\n{texto}",
            max_tokens=512
        )
        return f"🌐 Traducción ({idioma_destino}):\n{respuesta.strip()}"
    except Exception:
        return "No pude realizar la traducción."


# ─── RESUMEN RÁPIDO DE TEXTO/LINK ──────────────────────────────────────────
def resumir_texto(texto):
    """Genera un resumen conciso de un texto largo."""
    if len(texto) < 100:
        return "El texto es demasiado corto para resumir."
    try:
        respuesta = preguntar_ia(
            f"Resume el siguiente texto en 3-5 oraciones claras y concisas:\n\n{texto[:4000]}",
            max_tokens=300
        )
        return f"📄 Resumen:\n{respuesta.strip()}"
    except Exception:
        return "No pude generar el resumen."


# ─── GENERADOR DE CONTENIDO IA MEJORADO ─────────────────────────────────────
def generar_contenido_ia(instrucciones, tipo="respuesta"):
    """Genera contenido estructurado usando IA con diferentes modos."""
    try:
        if tipo == "lista":
            prompt = f"Genera una lista organizada basada en: {instrucciones}\nResponde con viñetas numeradas."
        elif tipo == "puntos":
            prompt = f"Extrae los puntos clave de: {instrucciones}\nResponde solo con los puntos, uno por línea."
        elif tipo == "codigo":
            prompt = f"Genera código para: {instrucciones}\nIncluye comentarios y manejos de errores."
        else:
            prompt = instrucciones
        respuesta = preguntar_ia(prompt, max_tokens=2048)
        return respuesta
    except Exception as e:
        return f"Error al generar contenido: {e}"


def procesar_mensaje(voz, silencio=False):
    def responder(texto):
        agregar_mensaje("JARVIS", texto)
        if not silencio:
            hablar(texto)
        historial_chat.append({"role": "assistant", "content": texto})
        guardar_historial()

    if not silencio:
        if "\n\nContenido del archivo" not in voz:
            agregar_mensaje("TÚ", voz)

    if "\n\nContenido del archivo" in voz:
        set_estado("PROCESANDO...")
        reproducir_sfx("processing")
        responder_streaming(f"Eres J.A.R.V.I.S, asistente formal de Pedro. Responde directamente. {voz}", silencio)
        return

    voz_lower = voz.lower()    # ─── COMANDOS IA AVANZADOS v2.0 ────────────────────────────────────────

    # 📋 ESCANEAR DOCUMENTO
    if any(x in voz_lower for x in ["escanea el documento", "analiza el documento", "escanear documento", "analizar documento", "qué dice este archivo"]):
        def bg_escanear():
            res = escanear_documento_ia()
            responder(res)
        threading.Thread(target=bg_escanear, daemon=True).start()

    # 📄 CREAR ARCHIVO CON IA
    elif any(x in voz_lower for x in ["crear archivo", "crea un archivo", "genera un archivo", "generar archivo"]):
        instruccion = voz
        for kw in ["crear archivo:", "crea un archivo:", "genera un archivo:", "generar archivo:"]:
            if kw in voz_lower:
                instruccion = voz.split(kw)[-1].strip()
                break
        else:
            for kw in ["crear archivo", "crea un archivo", "genera un archivo", "generar archivo"]:
                instruccion = instruccion.replace(kw, "").strip()
        if instruccion:
            def bg_crear_archivo():
                res = crear_archivo_ia(instruccion)
                responder(res)
            threading.Thread(target=bg_crear_archivo, daemon=True).start()
        else:
            responder("¿Qué tipo de archivo quieres que cree y qué debe contener?")

    # 💾 GUARDAR ARCHIVO EDITADO
    elif any(x in voz_lower for x in ["guardar archivo", "guarda el archivo", "guardar cambios", "guarda cambios"]):
        if editor_frame.winfo_ismapped():
            guardar_editor()
            responder("Archivo guardado correctamente.")
        elif archivo_cargado.get("ruta_original"):
            guardar_editor()
            responder("Archivo guardado correctamente.")
        else:
            responder("No hay ningún archivo abierto para guardar.")

    # 💾 GUARDAR COMO (en Documentos)
    elif any(x in voz_lower for x in ["guardar como", "guardar en documentos", "guardar en la carpeta"]):
        if editor_frame.winfo_ismapped():
            contenido = editor_text_box.get("1.0", tk.END).rstrip("\n")
            nombre = archivo_cargado.get("nombre", "archivo.txt")
            ruta_nueva = f"C:/Jarvis/Documentos/{nombre}"
            if guardar_archivo_nuevo(ruta_nueva, contenido):
                archivo_cargado["ruta_original"] = ruta_nueva
                responder(f"✅ Archivo guardado en Documentos: {nombre}")
            else:
                responder("No pude guardar el archivo en Documentos.")
        else:
            responder("No hay editor abierto. Carga un archivo primero.")

    # ✏️ MODIFICAR ARCHIVO CON IA
    elif any(x in voz_lower for x in ["modificar archivo", "modifica el archivo", "edita el archivo", "editar archivo"]):
        if archivo_cargado.get("nombre"):
            instruccion_mod = voz
            for kw in ["modificar archivo:", "modifica el archivo:", "edita el archivo:", "editar archivo:"]:
                if kw in voz_lower:
                    instruccion_mod = voz.split(kw)[-1].strip()
                    break
            if not instruccion_mod:
                responder("¿Qué cambios debo hacer en el archivo?")
            else:
                def bg_modificar():
                    ruta = f"C:/Jarvis/Documentos/{archivo_cargado['nombre']}"
                    if not os.path.exists(ruta):
                        ruta = f"C:/Jarvis/{archivo_cargado['nombre']}"
                    res = modificar_archivo_ia(ruta, instruccion_mod)
                    responder(res)
                threading.Thread(target=bg_modificar, daemon=True).start()
        else:
            responder("Primero carga un archivo con el botón + para poder modificarlo.")

    # 📖 LEER ARCHIVO POR VOZ
    elif any(x in voz_lower for x in ["lee el archivo", "leer archivo", "lee", "abre y lee"]):
        nombre_arch = voz_lower
        for kw in ["lee el archivo", "leer archivo", "lee", "abre y lee"]:
            nombre_arch = nombre_arch.replace(kw, "").strip()
        if nombre_arch:
            res = leer_archivo_por_voz(nombre_arch)
            responder(res)
        elif archivo_cargado.get("contenido"):
            contenido = archivo_cargado["contenido"][:2000]
            responder(f"📄 {archivo_cargado['nombre']}:\n\n{contenido}")
        else:
            responder("¿Qué archivo debo leer?")

    # 🔍 ANALIZAR CÓDIGO
    elif any(x in voz_lower for x in ["analiza el código", "analizar código", "revisa el código", "revisar código"]):
        def bg_analizar():
            res = analizar_codigo_ia()
            responder(res)
        threading.Thread(target=bg_analizar, daemon=True).start()

    # ▶️ EJECUTAR CÓDIGO PYTHON
    elif any(x in voz_lower for x in ["ejecuta el código", "ejecutar código", "corre el código", "run code"]):
        if archivo_cargado.get("contenido") and archivo_cargado.get("nombre", "").endswith(".py"):
            def bg_ejecutar():
                res = ejecutar_codigo_python(archivo_cargado["contenido"])
                responder(res)
            threading.Thread(target=bg_ejecutar, daemon=True).start()
        else:
            responder("Carga un archivo .py primero para poder ejecutarlo.")

    # 🔍 BUSCAR ARCHIVOS
    elif any(x in voz_lower for x in ["busca el archivo", "buscar archivo", "encuentra el archivo"]):
        nombre_buscar = voz_lower
        for kw in ["busca el archivo", "buscar archivo", "encuentra el archivo"]:
            nombre_buscar = nombre_buscar.replace(kw, "").strip()
        if nombre_buscar:
            res = buscar_archivos(nombre_buscar)
            responder(res)
        else:
            responder("¿Qué archivo debo buscar?")

    # 📁 LISTAR CARPETA
    elif any(x in voz_lower for x in ["qué hay en", "lista la carpeta", "contenidos de", "qué archivos hay"]):
        ruta_lista = "C:/Jarvis"
        for kw in ["qué hay en", "lista la carpeta", "contenidos de", "qué archivos hay"]:
            if kw in voz_lower:
                ruta_extraida = voz_lower.split(kw)[-1].strip()
                if ruta_extraida:
                    ruta_lista = f"C:/Jarvis/{ruta_extraida}"
                break
        res = listar_carpeta(ruta_lista)
        responder(res)

    # 🌐 TRADUCCIÓN
    elif any(x in voz_lower for x in ["traduce", "traducir", "cómo se dice", "cómo se escribe"]):
        prompt_ia = f"Traduce al español este texto: {voz}"
        def bg_traducir():
            res = preguntar_ia(prompt_ia, max_tokens=512)
            responder(f"🌐 Traducción:\n{res}")
        threading.Thread(target=bg_traducir, daemon=True).start()

    # 📋 COPIAR ÚLTIMA RESPUESTA
    elif any(x in voz_lower for x in ["copia eso", "copiar respuesta", "cópialo"]):
        if historial_chat:
            ultimas = [m for m in historial_chat if m["role"] == "assistant"]
            if ultimas:
                ventana.clipboard_clear()
                ventana.clipboard_append(ultimas[-1]["content"])
                responder("Copiado al portapapeles, señor.")
            else:
                responder("No hay respuestas que copiar.")
        else:
            responder("No hay respuestas que copiar.")

    # ─── COMANDOS ORIGINALES ────────────────────────────────────────────────

    # 📧 CORREO
    elif any(x in voz_lower for x in ["redactar correo", "redactar gmail", "escribe un correo"]):
        def bg_gmail():
            res = crear_gmail(voz)
            responder(res)
        threading.Thread(target=bg_gmail, daemon=True).start()

    # 🚀 ABRIR PROGRAMAS
    elif any(x in voz_lower for x in ["abrir programa", "ejecutar"]):
        res = abrir_programa_o_ruta(voz)
        responder(res)

    # 🔊 AUDIO
    elif any(x in voz_lower for x in ["salida de audio", "salida de sonido", "dispositivo de audio", "panel de sonido"]):
        nombre_pedido = _extraer_nombre_salida_pedida(_normalizar_texto(voz))
        if nombre_pedido:
            def bg_audio():
                flujo_cambiar_salida_audio(nombre_pedido)
            threading.Thread(target=bg_audio, daemon=True).start()
        else:
            res = cambiar_salida_audio()
            responder(res)

    # 🌤️ CLIMA
    elif any(x in voz_lower for x in ["clima", "tiempo hace"]):
        res = reporte_clima()
        responder(res)

    # 🎮 GAMING
    elif any(x in voz_lower for x in ["modo gaming", "modo juego"]):
        res = modo_gaming()
        responder(res)

    # 📧 ABRIR CORREO
    elif any(x in voz_lower for x in ["abrir correo", "abre el correo", "abrir gmail", "correo"]):
        webbrowser.open("https://mail.google.com")
        responder("Abriendo tu bandeja de correo principal, señor.")

    # 📋 TAREAS
    elif any(x in voz_lower for x in ["abrir tareas", "abre las tareas", "tareas"]):
        webbrowser.open("https://tasks.google.com")
        responder("Abriendo tus tareas pendientes.")

    # 🏫 PLATAFORMA / CLASSROOM
    elif any(x in voz_lower for x in ["plataforma", "abrir plataforma", "campus"]):
        responder("Abriendo la plataforma e iniciando sesión.")
        threading.Thread(target=abrir_plataforma_con_login, daemon=True).start()
    elif any(x in voz_lower for x in ["clases", "classroom"]):
        webbrowser.open("https://classroom.google.com/h")
        responder("Abriendo Google Classroom.")

    # 🔴 CERRAR PROCESOS
    elif any(x in voz_lower for x in ["cierra chrome", "cerrar chrome"]):
        subprocess.run(["taskkill", "/f", "/im", "chrome.exe"], creationflags=NO_WINDOW)
        mostrar_toast("Chrome cerrado", "success")
        responder("Cerrando proceso Chrome.")
    elif any(x in voz_lower for x in ["cierra opera gx", "cerrar opera gx"]):
        subprocess.run(["taskkill", "/f", "/im", "opera.exe"], creationflags=NO_WINDOW)
        mostrar_toast("Opera GX cerrado", "success")
        responder("Cerrando proceso Opera GX.")

    # 🌐 NAVEGADORES
    elif any(x in voz_lower for x in ["abre chrome", "abrir chrome"]):
        os.startfile(r"C:\Program Files\Google\Chrome\Application\chrome.exe")
        responder("Abriendo Google Chrome.")
    elif any(x in voz_lower for x in ["abre opera", "abrir opera", "opera gx"]):
        os.startfile(r"C:\Users\MSi\AppData\Local\Programs\Opera GX\opera.exe")
        responder("Abriendo navegador Opera GX.")

    # ⏰ RECORDATORIOS
    elif any(x in voz_lower for x in ["ponme un recordatorio", "recuérdame"]):
        hora_match = re.search(r'(\d{1,2})[:\s\.](\d{2})', voz)
        if hora_match:
            hora_str = f"{hora_match.group(1)}:{hora_match.group(2)}"
            for kw in ["recuérdame", "ponme un recordatorio"]:
                if kw in voz_lower:
                    mensaje_rec = voz_lower.split(kw)[-1].strip()
                    break
            else:
                mensaje_rec = voz
            hora_final = recordatorio(mensaje_rec, hora_str)
            if hora_final:
                responder(f"Recordatorio programado para las {hora_final}.")
        else:
            responder("No procesé la hora. Especifica en formato de reloj, como 'a las 6:30'.")

    # 🎵 MÚSICA
    elif any(x in voz_lower for x in ["reproduce", "escuchar", "reproducir"]):
        nombre = voz_lower
        for x in ["reproduce", "pon", "escuchar", "reproducir", "la cancion"]:
            nombre = nombre.replace(x, "").strip()
        reproducir_cancion(nombre, silencio)

    # 🔊 VOLUMEN
    elif any(x in voz_lower for x in ["sube el volumen", "subir volumen"]):
        nuevo = ajustar_volumen("subir")
        responder(f"Sistemas de audio incrementados al {nuevo}%.")
    elif any(x in voz_lower for x in ["baja el volumen", "bajar volumen"]):
        nuevo = ajustar_volumen("bajar")
        responder(f"Sistemas de audio reducidos al {nuevo}%.")
    elif any(x in voz_lower for x in ["pone volumen al", "volumen a"]):
        m_vol = re.search(r'(\d+)', voz_lower)
        if m_vol:
            nuevo = ajustar_volumen("set", int(m_vol.group(1)))
            responder(f"Volumen ajustado al {nuevo}%.")
        else:
            responder("¿A qué porcentaje debo poner el volumen?")

    # 📄 CREAR DOCUMENTOS
    elif any(x in voz_lower for x in ["crea un word", "crear word", "documento word"]):
        threading.Thread(target=lambda: crear_word(voz), daemon=True).start()
    elif any(x in voz_lower for x in ["crea un powerpoint", "crear powerpoint", "powerpoint"]):
        threading.Thread(target=lambda: crear_pptx(voz), daemon=True).start()
    elif any(x in voz_lower for x in ["crea un excel", "crear excel", "excel"]):
        threading.Thread(target=lambda: crear_excel(voz), daemon=True).start()
    elif any(x in voz_lower for x in ["crear python", "código python", "script python"]):
        threading.Thread(target=lambda: crear_archivo_py(voz), daemon=True).start()

    # 🔍 BUSQUEDA WEB
    elif any(x in voz_lower for x in ["busca en google", "buscar en google", "google"]):
        consulta = voz_lower
        for kw in ["busca en google", "buscar en google", "google"]:
            consulta = consulta.replace(kw, "")
        consulta = consulta.strip()
        if consulta:
            url = f"https://www.google.com/search?q={urllib.parse.quote(consulta)}"
            webbrowser.open(url)
            responder(f"Buscando '{consulta}' en Google, señor.")
        else:
            webbrowser.open("https://www.google.com")
            responder("Abriendo Google.")

    # 🎙️ HABLAR MÁS ALTO / BAJO
    elif any(x in voz_lower for x in ["habla más fuerte", "habla fuerte", "sube la voz"]):
        global voz_jarvis
        voz_jarvis = "es-MX-JorgeNeural"  # fallback cuando OpenAI TTS no está disponible  # Reset to default voice
        responder("Volumen de voz ajustado, señor.")

    # 📊 RESUMEN DEL DÍA
    elif any(x in voz_lower for x in ["resumen del día", "qué hice hoy", "historial de hoy"]):
        notas_hoy = []
        for nota in cargar_notas():
            if datetime.now().strftime("%Y-%m-%d") in nota.get("fecha", ""):
                notas_hoy.append(nota)
        historial_hoy = [m for m in historial_chat[-30:] if m.get("role") == "user"]
        resumen = f"📊 Resumen del día {datetime.now().strftime('%d/%m/%Y')}:\n"
        resumen += f"   • Mensajes en conversación: {len(historial_hoy)}\n"
        resumen += f"   • Notas creadas hoy: {len(notas_hoy)}\n"
        if notas_hoy:
            for n in notas_hoy:
                resumen += f"     → {n.get('titulo', 'Sin título')}\n"
        responder(resumen)

    # 🔄 LIMPIAR HISTORIAL
    elif any(x in voz_lower for x in ["limpiar historial", "borrar historial", "limpia el chat"]):
        historial_chat.clear()
        guardar_historial()
        for box in [chat_box, voz_chat_box]:
            box.config(state=tk.NORMAL)
            box.delete("1.0", tk.END)
            box.config(state=tk.DISABLED)
        try: chat_placeholder.place(relx=0.5, rely=0.45, anchor="center")
        except: pass
        try: voz_chat_placeholder.place(relx=0.5, rely=0.45, anchor="center")
        except: pass
        mostrar_toast("Historial limpiado", "success")
        responder("Historial limpiado, señor. Empezamos de cero.")

    # 🛑 APAGAR
    elif "apágate" in voz_lower or "hasta luego" in voz_lower:
        responder("Sistemas apagándose. Hasta luego, señor.")
        ventana.after(800, lambda: os._exit(0))

    else:
        # Llamada normal al Modelo LLM con Streaming
        set_estado("PROCESANDO...")
        reproducir_sfx("processing")
        responder_streaming(voz, silencio)


# ─── PALETA PROFESIONAL "OBSIDIAN & NEON" ──────────────────────────────────
BG          = "#03060D"      # Fondo ultra profundo
SIDEBAR_BG  = "#060A12"      # Sidebar con profundidad
SURFACE     = "#0A1020"      # Superficie de chats
SURFACE2    = "#101C35"      # Hover / tarjetas seleccionadas
SURFACE3    = "#182E50"      # Resaltes activos
ACCENT      = "#00CFFF"      # Cyan neón (identidad JARVIS)
ACCENT2     = "#2563EB"      # Azul eléctrico brillante
ACCENT_SOFT = "#0C2540"      # Cyan translúcido
GLOW        = "#7FEEFF"      # Blanco cyan de destello
TEXT        = "#E8F4FF"      # Texto claro premium
TEXT_DIM    = "#3D5A80"      # Texto secundario más sutil
GREEN       = "#10B981"      # Verde esmeralda
RED         = "#EF4444"      # Rojo error suave
AMBER       = "#F59E0B"      # Ámbar warning cálido
BORDER      = "#0E1A30"      # Borde sutil
BUBBLE_USER = "#0A2E52"      # Burbuja usuario (azul profundo)
BUBBLE_JAR  = "#081018"      # Burbuja JARVIS (obsidiana pura)

config_path = "C:/Jarvis/config.json"

def cargar_config():
    if os.path.exists(config_path):
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception: return {}
    return {}

def guardar_config(cfg):
    try:
        os.makedirs(os.path.dirname(config_path), exist_ok=True)
        with open(config_path, "w", encoding="utf-8") as f:
            json.dump(cfg, f, ensure_ascii=False, indent=2)
    except Exception: pass

_config_usuario = cargar_config()
if _config_usuario.get("color_burbuja_usuario"):
    BUBBLE_USER = _config_usuario["color_burbuja_usuario"]

FONT_TITLE  = ("Segoe UI Semibold", 15, "bold")
FONT_BODY   = ("Segoe UI", 11)
FONT_SMALL  = ("Segoe UI", 9)
FONT_MONO   = ("Cascadia Code", 10)
FONT_TIME   = ("Segoe UI", 10)
FONT_LOGO   = ("Segoe UI Semibold", 12, "bold")
FONT_MICRO  = ("Segoe UI", 8, "normal")

def _hex_a_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def _rgb_a_hex(rgb):
    return "#%02x%02x%02x" % tuple(int(max(0, min(255, c))) for c in rgb)

def _mezclar_color(c1, c2, t):
    r1, g1, b1 = _hex_a_rgb(c1)
    r2, g2, b2 = _hex_a_rgb(c2)
    return _rgb_a_hex((r1 + (r2-r1)*t, g1 + (g2-g1)*t, b1 + (b2-b1)*t))

def _animar_color(set_fn, color_inicio, color_fin, pasos=8, delay=14):
    def _paso(i=0):
        try:
            set_fn(_mezclar_color(color_inicio, color_fin, i / pasos))
        except tk.TclError: return
        if i < pasos:
            ventana.after(delay, lambda: _paso(i + 1))
    _paso()

# ─── VENTANA PRINCIPAL ────────────────────────────────────────────────────────
ventana = tk.Tk()
ventana.title("J.A.R.V.I.S. — Inteligencia Artificial")

# Icono custom: "J" sobre fondo oscuro (reemplaza la pluma de Tkinter)
_icono_img = tk.PhotoImage(width=32, height=32)
for _iy in range(32):
    for _ix in range(32):
        _dx, _dy = _ix - 16, _iy - 16
        _r2 = _dx*_dx + _dy*_dy
        if _r2 <= 256:  # círculo radio 16
            _dist = _r2 ** 0.5
            if _dist > 14:  # borde exterior
                _icono_img.put("#00CFFF", (_ix, _iy))
            elif _dist > 12:  # anillo
                _icono_img.put("#0A2744", (_ix, _iy))
            else:  # interior
                _icono_img.put("#060A12", (_ix, _iy))
        else:
            _icono_img.put("#060A12", (_ix, _iy))
# Dibujar "J" simple con píxeles
for _px, _py in [(12,7),(13,7),(14,7),(15,7),(16,7),
                 (15,8),(15,9),(15,10),(15,11),(15,12),
                 (14,13),(13,14),(12,15),(11,16),(10,17),(9,18),(9,19),(9,20),(10,21),(11,22),(12,23),
                 (12,24),(11,25),(10,26),(10,27)]:
    _icono_img.put("#00CFFF", (_px, _py))
try:
    ventana.iconphoto(True, _icono_img)
except Exception:
    pass

import queue
_ui_queue = queue.Queue()

# ─── MINI MODE: Esfera flotante en el escritorio ─────────────────────────────
_mini_mode = [True]  # arranca en mini mode
_mini_size = 140  # px del widget mini
_full_geo = [None]  # guardar geometria completa
_mini_layout_saved = [False]


def _set_frames_bg(color):
    """Cambia el bg de todos los frames ancestros del canvas para transparencia."""
    for w in (ventana, frame_main, content, panel_voz, voz_content):
        try: w.configure(bg=color)
        except: pass
    try: canvas.configure(bg=color, highlightthickness=0)
    except: pass


def abrir_mini_mode():
    """Oculta sidebar + chat, muestra solo la esfera en mini ventana transparente."""
    _mini_mode[0] = True
    ventana.update_idletasks()
    # Guardar geometria completa
    if not _mini_layout_saved[0]:
        _full_geo[0] = ventana.geometry()
        _mini_layout_saved[0] = True
    # Ocultar sidebar, accent y chat (solo canvas queda)
    sidebar.pack_forget()
    sidebar_accent.pack_forget()
    voz_chat_outer.pack_forget()
    # content sigue visible para el canvas
    content.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    panel_voz.pack(fill=tk.BOTH, expand=True)
    voz_content.pack(fill=tk.BOTH, expand=True, padx=10, pady=(10, 10))
    canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    # Fondo negro para transparencia
    _set_frames_bg("black")
    # Configurar ventana mini
    ventana.overrideredirect(True)
    ventana.attributes("-topmost", True)
    x_pos = 10
    y_pos = ventana.winfo_screenheight() - _mini_size - 50
    ventana.geometry(f"{_mini_size}x{_mini_size}+{x_pos}+{y_pos}")
    try:
        ventana.attributes("-transparentcolor", "black")
    except Exception:
        pass
    ventana.update_idletasks()
    ventana.deiconify()
    ventana.lift()
    set_estado("MINI MODE")


def abrir_ventana_completa():
    """Restaura la ventana completa con sidebar, chat y todo."""
    _mini_mode[0] = False
    # Quitar transparencia
    try:
        ventana.attributes("-transparentcolor", "")
    except Exception:
        pass
    ventana.overrideredirect(False)
    ventana.attributes("-topmost", False)
    # Restaurar geometria
    if _full_geo[0]:
        ventana.geometry(_full_geo[0])
    else:
        ventana.geometry("1100x700")
    # Restaurar colores de todos los frames
    _set_frames_bg(BG)
    # Restaurar sidebar (izquierda)
    sidebar.pack(side=tk.LEFT, fill=tk.Y, before=content)
    sidebar.pack_propagate(False)
    sidebar_accent.pack(side=tk.LEFT, fill=tk.Y, before=content)
    # Restaurar content con su layout interno
    content.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    panel_voz.pack(fill=tk.BOTH, expand=True)
    voz_content.pack(fill=tk.BOTH, expand=True, padx=10, pady=(10, 10))
    canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    # Restaurar voz_chat_outer (derecha del canvas, dentro de voz_content)
    voz_chat_outer.pack(side=tk.RIGHT, fill=tk.Y, padx=(6, 0))
    voz_chat_outer.pack_propagate(False)
    ventana.update_idletasks()
    ventana.deiconify()
    ventana.lift()
    ventana.focus_force()
    set_estado("ESCUCHANDO...")


def toggle_mini_mode():
    """Alterna entre mini mode y ventana completa."""
    if _mini_mode[0]:
        abrir_ventana_completa()
    else:
        abrir_mini_mode()

# ══════════════════════════════════════════════════════════════════════════════
# JCODEX — Agente de código AI integrado
# ══════════════════════════════════════════════════════════════════════════════
_jcodex_panel = None       # Frame del panel JCodex (se crea una vez)
_jcodex_active = [False]   # Está en modo JCodex?
_jcodex_folder = ["C:/Jarvis"]  # Carpeta del proyecto
_jcodex_chat_history = []  # Historial del agente

def _jc_crear_panel():
    """Crea el panel JCodex completo (file tree + editor + chat)."""
    global _jcodex_panel
    if _jcodex_panel:
        return
    _jcodex_panel = tk.Frame(content, bg=BG)

    # ── HEADER ──
    hdr = tk.Frame(_jcodex_panel, bg=SURFACE, pady=6)
    hdr.pack(fill=tk.X)
    tk.Label(hdr, text="  🖥️  JCodex — Agente de Código", font=("Segoe UI Semibold", 11, "bold"),
             fg="#00FFAA", bg=SURFACE).pack(side=tk.LEFT)
    tk.Label(hdr, text="Gemini-Powered", font=FONT_SMALL, fg=TEXT_DIM, bg=SURFACE).pack(side=tk.LEFT, padx=10)

    # Botón volver
    btn_back = tk.Label(hdr, text="  ✕ Cerrar  ", font=FONT_BODY, fg=RED, bg=SURFACE, cursor="hand2")
    btn_back.pack(side=tk.RIGHT, padx=6)
    btn_back.bind("<Button-1>", lambda e: _jc_volver())

    # Botón modo suspenso
    btn_suspenso = tk.Label(hdr, text="  ▶ Probar  ", font=FONT_BODY, fg=GREEN, bg=SURFACE, cursor="hand2")
    btn_suspenso.pack(side=tk.RIGHT, padx=6)
    btn_suspenso.bind("<Button-1>", lambda e: _jc_probar())

    tk.Frame(_jcodex_panel, bg=BORDER, height=1).pack(fill=tk.X)

    # ── BODY: file tree (izq) + editor/chat (der) ──
    body = tk.Frame(_jcodex_panel, bg=BG)
    body.pack(fill=tk.BOTH, expand=True)

    # ── FILE TREE (izquierda) ──
    tree_frame = tk.Frame(body, bg=SURFACE, width=220)
    tree_frame.pack(side=tk.LEFT, fill=tk.Y)
    tree_frame.pack_propagate(False)

    tk.Label(tree_frame, text=" 📁 Proyecto", font=("Segoe UI Semibold", 10, "bold"),
             fg=ACCENT, bg=SURFACE, anchor="w").pack(fill=tk.X, padx=8, pady=(8, 4))

    # Botón cambiar carpeta
    btn_folder = tk.Label(tree_frame, text="  📂 Cambiar carpeta...", font=FONT_SMALL,
                          fg=TEXT_DIM, bg=SURFACE, cursor="hand2", anchor="w")
    btn_folder.pack(fill=tk.X, padx=4)
    btn_folder.bind("<Button-1>", lambda e: _jc_cambiar_carpeta())

    tk.Frame(tree_frame, bg=BORDER, height=1).pack(fill=tk.X, padx=8, pady=4)

    # Scrollable file tree
    tree_canvas = tk.Canvas(tree_frame, bg=SURFACE, highlightthickness=0)
    tree_scroll = tk.Scrollbar(tree_frame, orient="vertical", command=tree_canvas.yview)
    tree_inner = tk.Frame(tree_canvas, bg=SURFACE)
    tree_inner.bind("<Configure>", lambda e: tree_canvas.configure(scrollregion=tree_canvas.bbox("all")))
    tree_canvas.create_window((0, 0), window=tree_inner, anchor="nw")
    tree_canvas.configure(yscrollcommand=tree_scroll.set)
    tree_scroll.pack(side=tk.RIGHT, fill=tk.Y)
    tree_canvas.pack(fill=tk.BOTH, expand=True)

    # ── RIGHT: editor + chat ──
    right = tk.Frame(body, bg=BG)
    right.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    # Editor (arriba)
    editor_frame = tk.Frame(right, bg=SURFACE)
    editor_frame.pack(fill=tk.BOTH, expand=True, padx=(1, 0), pady=(0, 1))
    tk.Label(editor_frame, text=" 📝 Editor", font=("Segoe UI Semibold", 10, "bold"),
             fg=ACCENT, bg=SURFACE, anchor="w").pack(fill=tk.X, padx=8, pady=(6, 2))
    editor_text = tk.Text(editor_frame, bg="#0A0E1A", fg=TEXT, font=("Consolas", 10),
                          insertbackground=ACCENT, selectbackground="#1a3a5c",
                          wrap=tk.WORD, relief=tk.FLAT, padx=8, pady=4,
                          state=tk.DISABLED, undo=True)
    editor_scroll = tk.Scrollbar(editor_frame, command=editor_text.yview)
    editor_text.configure(yscrollcommand=editor_scroll.set)
    editor_scroll.pack(side=tk.RIGHT, fill=tk.Y)
    editor_text.pack(fill=tk.BOTH, expand=True)

    tk.Frame(right, bg=BORDER, height=1).pack(fill=tk.X)

    # Chat del agente (abajo)
    chat_frame = tk.Frame(right, bg=BG, height=180)
    chat_frame.pack(fill=tk.BOTH, padx=(1, 0))
    chat_frame.pack_propagate(False)
    tk.Label(chat_frame, text=" 🤖 Chat con JCodex", font=("Segoe UI Semibold", 10, "bold"),
             fg="#00FFAA", bg=BG, anchor="w").pack(fill=tk.X, padx=8, pady=(6, 2))
    chat_box = tk.Text(chat_frame, bg="#0A0E1A", fg=TEXT, font=("Consolas", 10),
                       insertbackground=ACCENT, wrap=tk.WORD, relief=tk.FLAT,
                       padx=8, pady=4, state=tk.DISABLED, height=6)
    chat_scroll = tk.Scrollbar(chat_frame, command=chat_box.yview)
    chat_box.configure(yscrollcommand=chat_scroll.set)
    chat_scroll.pack(side=tk.RIGHT, fill=tk.Y)
    chat_box.pack(fill=tk.BOTH, expand=True)

    # Input bar
    input_bar = tk.Frame(chat_frame, bg=SURFACE)
    input_bar.pack(fill=tk.X, padx=0, pady=(2, 0))
    chat_entry = tk.Entry(input_bar, bg="#0A0E1A", fg=TEXT, font=("Consolas", 10),
                          insertbackground=ACCENT, relief=tk.FLAT, highlightthickness=1,
                          highlightcolor=ACCENT, highlightbackground=BORDER)
    chat_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=4, pady=4, ipady=4)
    btn_send = tk.Label(input_bar, text="  ➤  ", font=("Segoe UI", 12, "bold"),
                        fg="#00FFAA", bg=SURFACE, cursor="hand2")
    btn_send.pack(side=tk.RIGHT, padx=4, pady=4)
    chat_entry.bind("<Return>", lambda e: _jc_enviar_mensaje())
    btn_send.bind("<Button-1>", lambda e: _jc_enviar_mensaje())

    # Guardar referencias para acceso rápido
    _jcodex_panel._tree_inner = tree_inner
    _jcodex_panel._tree_canvas = tree_canvas
    _jcodex_panel._editor = editor_text
    _jcodex_panel._chat_box = chat_box
    _jcodex_panel._chat_entry = chat_entry
    _jcodex_panel._file_labels = []
    _jcodex_panel._current_file = [None]
    _jcodex_panel._modified_files = {}  # {path: original_content}

def _jc_mostrar():
    """Muestra el panel JCodex ocultando el contenido normal."""
    global mic_activo
    _jc_crear_panel()
    _jcodex_active[0] = True
    # Pausar micrófono (JCodex solo usa texto)
    if mic_activo:
        mic_activo = False
        try:
            sidebar_btn_refs["mic_lbl"].config(text="Micrófono OFF", fg=RED)
            sidebar_btn_refs["mic_ico"].config(fg=RED)
        except: pass
        set_estado("MODO JCODEX — TEXTO")
    # Ocultar widgets normales
    try:
        voz_content.pack_forget()
    except: pass
    try:
        panel_voz.pack_forget()
    except: pass
    _jcodex_panel.pack(fill=tk.BOTH, expand=True)
    _jc_refrescar_arbol()
    _jc_chat_mensaje("system", "JCodex listo. Selecciona una carpeta y hazme una solicitud.")

def _jc_volver():
    """Vuelve al modo normal desde JCodex."""
    global _jcodex_active, mic_activo
    _jcodex_active[0] = False
    if _jcodex_panel:
        _jcodex_panel.pack_forget()
    panel_voz.pack(fill=tk.BOTH, expand=True)
    voz_content.pack(fill=tk.BOTH, expand=True, padx=10, pady=(10, 10))
    canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    voz_chat_outer.pack(side=tk.RIGHT, fill=tk.Y, padx=(6, 0))
    voz_chat_outer.pack_propagate(False)
    # Reactivar micrófono al volver
    mic_activo = True
    try:
        sidebar_btn_refs["mic_lbl"].config(text="Micrófono ON", fg=GREEN)
        sidebar_btn_refs["mic_ico"].config(fg=GREEN)
    except: pass
    set_estado("ESCUCHANDO...")

def _jc_cambiar_carpeta():
    """Abre selector de carpeta."""
    folder = fd.askdirectory(initialdir=_jcodex_folder[0], title="Seleccionar carpeta del proyecto")
    if folder:
        _jcodex_folder[0] = folder
        _jc_refrescar_arbol()
        _jc_chat_mensaje("system", f"Carpeta: {folder}")

def _jc_refrescar_arbol():
    """Refresca el árbol de archivos."""
    if not _jcodex_panel: return
    tree_inner = _jcodex_panel._tree_inner
    # Limpiar
    for w in tree_inner.winfo_children():
        w.destroy()
    _jcodex_panel._file_labels = []
    # Llenar
    folder = _jcodex_folder[0]
    if not os.path.isdir(folder): return
    for item in sorted(os.listdir(folder)):
        full = os.path.join(folder, item)
        if os.path.isdir(full): continue  # solo archivos por ahora
        if item.endswith((".py", ".js", ".ts", ".html", ".css", ".json", ".md", ".txt", ".env", ".vbs")):
            lbl = tk.Label(tree_inner, text=f"  📄 {item}", font=FONT_SMALL, fg=TEXT_DIM,
                           bg=SURFACE, anchor="w", cursor="hand2")
            lbl.pack(fill=tk.X, padx=4, pady=1)
            lbl.bind("<Button-1>", lambda e, p=full: _jc_abrir_archivo(p))
            _jcodex_panel._file_labels.append(lbl)

def _jc_abrir_archivo(path):
    """Abre un archivo en el editor."""
    if not _jcodex_panel: return
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
    except Exception as e:
        _jc_chat_mensaje("system", f"Error al leer: {e}")
        return
    editor = _jcodex_panel._editor
    editor.config(state=tk.NORMAL)
    editor.delete("1.0", tk.END)
    editor.insert("1.0", content)
    editor.config(state=tk.DISABLED)
    _jcodex_panel._current_file[0] = path
    # Highlight nombre del archivo abierto
    for lbl in _jcodex_panel._file_labels:
        if path in (lbl.cget("text") or ""):
            lbl.config(fg="#00FFAA")
        else:
            lbl.config(fg=TEXT_DIM)
    _jc_chat_mensaje("system", f"Abierto: {os.path.basename(path)} ({len(content)} chars)")

def _jc_chat_mensaje(role, text):
    """Agrega un mensaje al chat del agente."""
    if not _jcodex_panel: return
    box = _jcodex_panel._chat_box
    box.config(state=tk.NORMAL)
    prefix = "🤖 " if role == "agent" else "📁 " if role == "user" else "⚙ "
    box.insert(tk.END, f"\n{prefix}{text}")
    box.config(state=tk.DISABLED)
    box.see(tk.END)

def _jc_enviar_mensaje():
    """Envía un mensaje al agente de código."""
    if not _jcodex_panel: return
    entry = _jcodex_panel._chat_entry
    msg = entry.get().strip()
    if not msg: return
    entry.delete(0, tk.END)
    _jc_chat_mensaje("user", msg)
    _jcodex_chat_history.append({"role": "user", "content": msg})
    # Ejecutar en hilo separado
    threading.Thread(target=_jc_procesar_solicitud, args=(msg,), daemon=True).start()

def _jc_llamar_ia_solicitud(system_prompt, messages, max_tokens=8192):
    """Llama a IA con fallback: Gemini (esperar 20s si 429) → NVIDIA NIM (sin espera)."""
    import time as _time
    # ── 1. Gemini: intentar 1 vez, si 429 esperar 20s y reintentar 1 vez más ──
    if gemini_client:
        for attempt in range(2):  # solo 2 intentos
            try:
                response = gemini_client.models.generate_content(
                    model=GEMINI_MODEL,
                    contents=[m["content"] for m in messages if m["role"] != "system"],
                    config=genai.types.GenerateContentConfig(
                        system_instruction=system_prompt,
                        max_output_tokens=max_tokens
                    )
                )
                return response.text, "Gemini"
            except Exception as e:
                err_str = str(e)
                if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
                    if attempt == 0:
                        _jc_chat_mensaje("system", "⏳ Cuota Gemini agotada. Esperando 20s para reintentar...")
                        _time.sleep(20)
                        continue
                    else:
                        _jc_chat_mensaje("system", "⚠️ Gemini agotado. Usando NVIDIA NIM...")
                else:
                    _log_error("_jc_llamar_ia_solicitud - Gemini", e)
                    break
    # ── 2. NVIDIA NIM: sin espera, directo ──
    if nvidia_client:
        try:
            nvidia_msgs = [{"role": "system", "content": system_prompt}] + [{"role": m["role"], "content": m["content"]} for m in messages if m["role"] != "system"]
            respuesta = nvidia_client.chat.completions.create(
                model=NVIDIA_MODEL,
                messages=nvidia_msgs,
                max_tokens=max_tokens
            )
            return respuesta.choices[0].message.content, "NVIDIA NIM"
        except Exception as e:
            _log_error("_jc_llamar_ia_solicitud - NIM", e)
            _jc_chat_mensaje("system", f"❌ NVIDIA NIM también falló: {e}")
    return None, None

def _jc_procesar_solicitud(msg):
    """Procesa una solicitud del usuario usando IA como agente de código."""
    if not gemini_client and not nvidia_client:
        _jc_chat_mensaje("system", "❌ No hay proveedores de IA disponibles. Verifica tus API keys.")
        return
    # Leer archivo actual si hay uno
    archivo_actual = ""
    archivo_path = ""
    if _jcodex_panel and _jcodex_panel._current_file[0]:
        archivo_path = _jcodex_panel._current_file[0]
        try:
            with open(archivo_path, "r", encoding="utf-8", errors="replace") as f:
                archivo_actual = f.read()
        except: pass
    # Construir contexto
    system_prompt = (
        "Eres JCodex, un agente de programación profesional integrado en JARVIS. "
        "Tienes conocimiento avanzado de: Python, JavaScript, TypeScript, Rust, Go, Java, C/C++, SQL, HTML/CSS, React, Node.js, Docker, Kubernetes, Git, Linux, y APIs REST/GraphQL. "
        "Puedes: leer archivos, crear archivos, modificarlos, ejecutar comandos del sistema, depurar errores, refactorizar código, escribir tests, y revisar PRs. "
        "Siempre escribe código limpio, bien documentado y siguiendo las mejores prácticas de la industria. "
        "Si el usuario pregunta sobre NVIDIA/GPU/accelerated computing, usa las NVIDIA Skills disponibles en nvidia_skills/. "
        "IMPORTANTE: Cuando modifiques un archivo, devuelve SOLO el código completo modificado usando este formato EXACTO:"
        "\n\n### EDITAR: ruta/archivo.ext\n```\n(código completo aquí)\n```\n\n"
        "Para crear un archivo nuevo: ### CREAR: ruta/nuevo_archivo.ext\n```\n(código completo)\n```\n\n"
        "Cuando necesites ejecutar un comando, usa: ### EJECUTAR: comando\n\n"
        "Si solo es una respuesta conversacional, responde normalmente. "
        f"Carpeta del proyecto: {_jcodex_folder[0]}\n"
        f"Archivo abierto: {archivo_path}\n"
    )
    if archivo_actual:
        system_prompt += f"\nContenido del archivo abierto:\n```\n{archivo_actual[:8000]}\n```"

    # ── Cargar NVIDIA Skills si el prompt menciona GPU/accelerated computing ──
    nvidia_skill = _detectar_nvidia_skill(msg)
    if nvidia_skill:
        system_prompt += f"\n\n--- NVIDIA SKILL (conocimiento técnico) ---\n{nvidia_skill[:6000]}\n--- FIN SKILL ---"
        _jc_chat_mensaje("system", "📚 NVIDIA Skill cargada para esta consulta")

    messages = [{"role": "system", "content": system_prompt}] + _jcodex_chat_history[-15:]
    respuesta, proveedor = _jc_llamar_ia_solicitud(system_prompt, messages)
    if not respuesta:
        _jc_chat_mensaje("system", "❌ No se pudo obtener respuesta de ningún proveedor.")
        return
    _jc_chat_mensaje("system", f"(vía {proveedor})")

    _jcodex_chat_history.append({"role": "assistant", "content": respuesta})
    _jc_chat_mensaje("agent", respuesta[:500])

    # ── Parsear acciones: EDITAR o EJECUTAR ──
    import re as _re
    # Buscar bloques de edición
    edits = _re.findall(r"### EDITAR:\s*(.+?)\n```(?:\w*)\n(.*?)```", respuesta, _re.DOTALL)
    for ruta, codigo in edits:
        ruta = ruta.strip()
        if not os.path.isabs(ruta):
            ruta = os.path.join(_jcodex_folder[0], ruta)
        # Guardar original para backup
        if os.path.exists(ruta) and ruta not in _jcodex_panel._modified_files:
            with open(ruta, "r", encoding="utf-8", errors="replace") as f:
                _jcodex_panel._modified_files[ruta] = f.read()
        try:
            with open(ruta, "w", encoding="utf-8") as f:
                f.write(codigo.strip())
            _jc_chat_mensaje("system", f"✅ Editado: {os.path.basename(ruta)}")
        except Exception as e:
            _jc_chat_mensaje("system", f"❌ Error al editar {ruta}: {e}")
    # Buscar bloques de creación de archivos nuevos
    creates = _re.findall(r"### CREAR:\s*(.+?)\n```(?:\w*)\n(.*?)```", respuesta, _re.DOTALL)
    for ruta, codigo in creates:
        ruta = ruta.strip()
        if not os.path.isabs(ruta):
            ruta = os.path.join(_jcodex_folder[0], ruta)
        try:
            os.makedirs(os.path.dirname(ruta), exist_ok=True)
            with open(ruta, "w", encoding="utf-8") as f:
                f.write(codigo.strip())
            _jc_chat_mensaje("system", f"✅ Creado: {os.path.basename(ruta)}")
        except Exception as e:
            _jc_chat_mensaje("system", f"❌ Error al crear {ruta}: {e}")
    # Buscar comandos
    cmds = _re.findall(r"### EJECUTAR:\s*(.+?)\n", respuesta)
    for cmd in cmds:
        cmd = cmd.strip()
        _jc_chat_mensaje("system", f"⚙ Ejecutando: {cmd}")
        try:
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=30,
                                    cwd=_jcodex_folder[0])
            output = result.stdout.strip() or result.stderr.strip() or "(sin salida)"
            _jc_chat_mensaje("system", output[:500])
        except Exception as e:
            _jc_chat_mensaje("system", f"❌ Error: {e}")
    # Refrescar árbol si hubo ediciones/creaciones
    if edits or creates:
        ui(_jc_refrescar_arbol)

def _jc_probar():
    """Crea jarvisP.py y entra en modo suspenso."""
    src = os.path.join(_jcodex_folder[0], "jarvis.py")
    dst = os.path.join(_jcodex_folder[0], "jarvisP.py")
    if not os.path.exists(src):
        _jc_chat_mensaje("system", "❌ No se encontró jarvis.py en la carpeta")
        return
    try:
        with open(src, "r", encoding="utf-8") as f:
            content = f.read()
        # Agregar indicador de modo suspenso al inicio
        with open(dst, "w", encoding="utf-8") as f:
            f.write("# ═══ MODO SUSPENSO — JARVIS DE PRUEBA (JCodex) ═══\n")
            f.write("# Este archivo se generó automáticamente para probar cambios.\n")
            f.write("# Al cerrarlo, JARVIS preguntará si aplicar los cambios.\n\n")
            f.write(content)
        _jc_chat_mensaje("system", "⏸️ Modo suspenso: jarvisP.py creado")
        _jc_chat_mensaje("system", "Ejecuta: python jarvisP.py")
        _jc_chat_mensaje("system", "Al cerrarlo, volveré a preguntar si aplicar los cambios.")
        # Lanzar jarvisP.py en proceso separado
        subprocess.Popen(["python", dst])
        # Monitorear cierre en hilo separado
        threading.Thread(target=_jc_monitorear_prueba, args=(dst, src), daemon=True).start()
    except Exception as e:
        _jc_chat_mensaje("system", f"❌ Error: {e}")

def _jc_monitorear_prueba(jarvisp_path, original_path):
    """Espera a que se cierre jarvisP.py y preguntar si aplicar."""
    # Esperar a que el proceso de jarvisP.py esté corriendo
    time.sleep(2)
    # Esperar a que termine
    while True:
        try:
            result = subprocess.run(
                ["tasklist", "/fi", "IMAGENAME eq python.exe", "/fo", "csv", "/nh"],
                capture_output=True, text=True, timeout=5
            )
            if jarvisp_path.replace("/", "\\") not in result.stdout:
                break
        except: pass
        time.sleep(2)
    # jarvisP.py terminó — preguntar
    ui(lambda: _jc_preguntar_aplicar(jarvisp_path, original_path))

def _jc_preguntar_aplicar(jarvisp_path, original_path):
    """Muestra diálogo preguntando si aplicar cambios."""
    from tkinter import messagebox
    respuesta = messagebox.askyesnocancel(
        "JCodex — Cambios en jarvis.py",
        "jarvisP.py terminó. ¿Deseas aplicar los cambios?\n\n"
        "✅ Sí → Reemplazar jarvis.py con jarvisP.py\n"
        "❌ No → Mantener jarvis.py original\n"
        "⬜ Cancelar → No hacer nada"
    )
    if respuesta is True:  # Sí → aplicar
        try:
            with open(jarvisp_path, "r", encoding="utf-8") as f:
                new_content = f.read()
            # Quitar cabecera de modo suspenso
            if new_content.startswith("# ═══ MODO SUSPENSO"):
                idx = new_content.find("\n", new_content.find("\n") + 1) + 1
                if new_content[idx:idx+2] == "\n\n":
                    idx += 2
                new_content = new_content[idx:]
            with open(original_path, "w", encoding="utf-8") as f:
                f.write(new_content)
            os.remove(jarvisp_path)
            _jc_chat_mensaje("system", "✅ jarvis.py actualizado con los cambios.")
            _jc_chat_mensaje("system", "Reinicia JARVIS para aplicar.")
        except Exception as e:
            _jc_chat_mensaje("system", f"❌ Error al aplicar: {e}")
    elif respuesta is False:  # No → mantener original
        if os.path.exists(jarvisp_path):
            os.remove(jarvisp_path)
        _jc_chat_mensaje("system", "📁 Cambios descartados. jarvis.py original mantenido.")
    else:  # Cancelar
        _jc_chat_mensaje("system", "⏸️ jarvisP.py conservado para revisión manual.")

def ui(func, *args, **kwargs):
    if threading.current_thread() is threading.main_thread():
        return func(*args, **kwargs)
    _ui_queue.put((func, args, kwargs))

def _procesar_cola_ui():
    while True:
        try:
            func, args, kwargs = _ui_queue.get_nowait()
        except queue.Empty: break
        try:
            func(*args, **kwargs)
        except Exception as e:
            _log_error("_procesar_cola_ui", e)
    ventana.after(30, _procesar_cola_ui)

ventana.configure(bg=BG)
ventana.geometry("1100x700")
ventana.minsize(950, 600)

frame_main = tk.Frame(ventana, bg=BG)
frame_main.pack(fill=tk.BOTH, expand=True)

# ─── SIDEBAR IZQUIERDO ────────────────────────────────────────────────────────
sidebar = tk.Frame(frame_main, bg=SIDEBAR_BG, width=196)
sidebar.pack(side=tk.LEFT, fill=tk.Y)
sidebar.pack_propagate(False)

# ─── LÍNEA ACENTO GRADIENTE EN EL BORDE DERECHO DEL SIDEBAR ───
sidebar_accent = tk.Frame(frame_main, bg=ACCENT, width=2)
sidebar_accent.pack(side=tk.LEFT, fill=tk.Y)

logo_canvas = tk.Canvas(sidebar, bg=SIDEBAR_BG, highlightthickness=0, width=150, height=150)
logo_canvas.pack(pady=(22, 4))

def _draw_jarvis_logo(c, glow=0.0):
    c.delete("all")
    cx, cy, R = 75, 75, 60
    halo_r = R + 8 + glow * 6
    halo_color = _mezclar_color(SIDEBAR_BG, ACCENT_SOFT, 0.25 + glow * 0.5)
    c.create_oval(cx-halo_r, cy-halo_r, cx+halo_r, cy+halo_r, outline=halo_color, width=1)
    
    for r, col in [(R, "#0d2a45"), (R-8, "#0f3a5c"), (R-16, "#1a5278")]:
        c.create_oval(cx-r, cy-r, cx+r, cy+r, outline=col, width=2, fill="")
        
    arco_color = _mezclar_color(ACCENT2, GLOW, glow * 0.6)
    c.create_arc(cx-R+2, cy-R+2, cx+R-2, cy+R-2, start=30, extent=60, outline=arco_color, width=2, style="arc")
    c.create_arc(cx-R+2, cy-R+2, cx+R-2, cy+R-2, start=210, extent=60, outline=arco_color, width=2, style="arc")
    c.create_arc(cx-R+9, cy-R+9, cx+R-9, cy+R-9, start=120, extent=25, outline=ACCENT_SOFT, width=1, style="arc")
    c.create_arc(cx-R+9, cy-R+9, cx+R-9, cy+R-9, start=300, extent=25, outline=ACCENT_SOFT, width=1, style="arc")
    
    c.create_oval(cx-38, cy-38, cx+38, cy+38, fill="#071520", outline="#1a4a6a", width=1)
    c.create_text(cx, cy-6, text="J.A.R.V.I.S", font=("Segoe UI", 9, "bold"), fill=ACCENT)
    c.create_text(cx, cy+10, text="◆", font=("Segoe UI", 8), fill=ACCENT2)

_draw_jarvis_logo(logo_canvas)

_logo_glow_fase = [0.0]
def _animar_logo_glow():
    _logo_glow_fase[0] += 0.045
    glow = (math.sin(_logo_glow_fase[0]) + 1) / 2
    _draw_jarvis_logo(logo_canvas, glow)
    ventana.after(60, _animar_logo_glow)
_animar_logo_glow()

hora_label = tk.Label(sidebar, text="", font=("Segoe UI", 9), fg=TEXT_DIM, bg=SIDEBAR_BG, anchor="center")
hora_label.pack(fill=tk.X, padx=16, pady=(2, 8))

status_pill = tk.Frame(sidebar, bg=SURFACE, padx=10, pady=6, highlightthickness=1, highlightbackground=BORDER)
status_pill.pack(padx=16, pady=(0, 14), fill=tk.X)

status_dot = tk.Canvas(status_pill, width=10, height=10, bg=SURFACE, highlightthickness=0)
status_dot.pack(side=tk.LEFT, padx=(0, 8))

estado_label = tk.Label(status_pill, text="INICIANDO...", font=("Segoe UI", 8, "bold"), fg=ACCENT, bg=SURFACE, anchor="w")
estado_label.pack(side=tk.LEFT, fill=tk.X, expand=True)

# ─── ESTADOS DE JARVIS (texto + color/animación de la bolita) ──────────────
# Cada entrada: (texto que se muestra, acción interna para animar la bolita)
_MAPA_ESTADOS_ACCION = [
    ("MIC DESACTIVADO",  "mic_off"),
    ("MODO TEXTO",       "modo_texto"),
    ("ESCUCHANDO",       "escuchando"),
    ("PENSANDO",         "pensando"),
    ("PROCESANDO",       "pensando"),
    ("PREPARANDO",       "pensando"),
    ("GENERANDO",        "pensando"),
]

def set_estado(texto):
    """Actualiza el texto de estado y la acción interna que usa la bolita
    para elegir color/animación. Es thread-safe (usa la cola de UI)."""
    global estado_accion
    accion = "pensando"
    for clave, acc in _MAPA_ESTADOS_ACCION:
        if clave in texto:
            accion = acc
            break
    estado_accion = accion
    ui(estado_label.config, text=texto)

def _color_estado_objetivo():
    """Color objetivo del estado actual (snap)."""
    if not mic_activo and not modo_texto: return RED
    if hablando: return AMBER
    return {
        "escuchando": GREEN,
        "pensando": ACCENT2,
        "mic_off": RED,
        "modo_texto": ACCENT,
    }.get(estado_accion, ACCENT)

_sphere_color_rgb = list(_hex_a_rgb(ACCENT))  # color actual suavizado
_sphere_color_target = list(_hex_a_rgb(ACCENT))
_sphere_color_prev = list(_hex_a_rgb(ACCENT))  # color del frame anterior

# ─── MORPH FACTOR (smooth deformation when speaking) ───
_morph_factor = 0.0     # 0.0 = esfera en reposo, 1.0 = esfera deformada (hablando)
_morph_velocity = 0.0   # velocidad del morph (para efecto elástico)

def _color_estado_actual():
    """Retorna el color suavizado (lerp) para la transición fluida.
    Usa easing suave con factor adaptativo para transiciones profesionales."""
    global _sphere_color_rgb, _sphere_color_target, _sphere_color_prev
    target = _hex_a_rgb(_color_estado_objetivo())
    _sphere_color_target = list(target)
    # Distancia actual al objetivo
    distancia = sum(abs(_sphere_color_target[i] - _sphere_color_rgb[i]) for i in range(3))
    # Factor adaptativo suave: nunca salta, siempre desliza
    if distancia > 100:
        factor = 0.18  # cambio grande → transición activa pero suave
    elif hablando:
        factor = 0.12  # hablando → reacciona al ritmo
    else:
        factor = 0.06  # idle → muy suave, elegante
    # Aplicar lerp con easing exponencial (suaviza el final)
    for i in range(3):
        diff = _sphere_color_target[i] - _sphere_color_rgb[i]
        _sphere_color_rgb[i] += diff * factor
    # Guardar prev para detectar cambios
    _sphere_color_prev = list(_sphere_color_rgb)
    return _rgb_a_hex(tuple(int(c) for c in _sphere_color_rgb))

_status_dot_fase = [0.0]
def _animar_status_dot():
    _status_dot_fase[0] += 0.18
    color = _color_estado_actual()
    pulso = (math.sin(_status_dot_fase[0]) + 1) / 2
    radio = 3 + pulso * 1.6
    status_dot.delete("all")
    status_dot.create_oval(5-radio, 5-radio, 5+radio, 5+radio, fill=color, outline="")
    ventana.after(90, _animar_status_dot)
_animar_status_dot()

tk.Frame(sidebar, bg=BORDER, height=1).pack(fill=tk.X, padx=16, pady=(0, 10))

sys_panel = tk.Frame(sidebar, bg=SIDEBAR_BG)
sys_panel.pack(fill=tk.X, padx=18, pady=(0, 4))

def _crear_barra_sistema(parent, etiqueta, color):
    fila = tk.Frame(parent, bg=SIDEBAR_BG)
    fila.pack(fill=tk.X, pady=5)
    lbl = tk.Label(fila, text=etiqueta, font=FONT_MICRO, fg=TEXT_DIM, bg=SIDEBAR_BG, anchor="w")
    lbl.pack(fill=tk.X)
    pista = tk.Frame(fila, bg=SURFACE2, height=4)
    pista.pack(fill=tk.X, pady=(3, 0))
    relleno = tk.Frame(pista, bg=color)
    relleno.place(x=0, y=0, relheight=1, relwidth=0.02)
    return lbl, relleno

cpu_lbl, cpu_fill = _crear_barra_sistema(sys_panel, "CPU   --%", ACCENT)
ram_lbl, ram_fill = _crear_barra_sistema(sys_panel, "RAM   --%", ACCENT2)

def _actualizar_panel_sistema():
    try:
        import psutil
        cpu = psutil.cpu_percent(interval=None)
        ram = psutil.virtual_memory().percent
        cpu_lbl.config(text=f"CPU   {cpu:>3.0f}%")
        ram_lbl.config(text=f"RAM   {ram:>3.0f}%")
        cpu_fill.place(relwidth=max(0.02, min(1.0, cpu / 100)))
        ram_fill.place(relwidth=max(0.02, min(1.0, ram / 100)))
    except Exception: pass
    ventana.after(2500, _actualizar_panel_sistema)
_actualizar_panel_sistema()

tk.Frame(sidebar, bg=BORDER, height=1).pack(fill=tk.X, padx=16, pady=(8, 4))

sidebar_btn_refs = {}

def make_sidebar_btn(parent, text, icon, color, cmd):
    f = tk.Frame(parent, bg=SIDEBAR_BG, cursor="hand2")
    f.pack(fill=tk.X, padx=8, pady=3)
    barra = tk.Frame(f, bg=SIDEBAR_BG, width=3)
    barra.pack(side=tk.LEFT, fill=tk.Y)
    ico = tk.Label(f, text=icon, font=("Segoe UI", 12), fg=color, bg=SIDEBAR_BG, width=3)
    ico.pack(side=tk.LEFT, padx=(6, 0), pady=7)
    lbl = tk.Label(f, text=text, font=FONT_BODY, fg=TEXT, bg=SIDEBAR_BG, anchor="w")
    lbl.pack(side=tk.LEFT, fill=tk.X, expand=True, pady=7)

    def _set_bg(col):
        for w in (f, ico, lbl): w.config(bg=col)

    def _on_enter(e):
        _animar_color(_set_bg, SIDEBAR_BG, SURFACE2, pasos=6, delay=12)
        _animar_color(lambda col: barra.config(bg=col), SIDEBAR_BG, ACCENT, pasos=6, delay=12)

    def _on_leave(e):
        _animar_color(_set_bg, SURFACE2, SIDEBAR_BG, pasos=6, delay=12)
        _animar_color(lambda col: barra.config(bg=col), ACCENT, SIDEBAR_BG, pasos=6, delay=12)

    def _click(e):
        reproducir_sfx("click")
        cmd()

    for w in (f, ico, lbl):
        w.bind("<Button-1>", _click)
        w.bind("<Enter>", _on_enter)
        w.bind("<Leave>", _on_leave)
    return ico, lbl

# ─── ÁREA DE CONTENIDO ────────────────────────────────────────────────────────
# bg se cambia a "black" en mini mode para transparencia
content = tk.Frame(frame_main, bg=BG)
content.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

# ─── PANEL VOZ ────────────────────────────────────────────────────────────────
panel_voz = tk.Frame(content, bg=BG)
panel_voz.pack(fill=tk.BOTH, expand=True)

voz_content = tk.Frame(panel_voz, bg=BG)
voz_content.pack(fill=tk.BOTH, expand=True, padx=10, pady=(10, 10))

canvas = tk.Canvas(voz_content, bg=BG, highlightthickness=0)
canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

def _on_canvas_click(e):
    """En mini mode, click en cualquier lugar abre la ventana completa."""
    if _mini_mode[0]:
        abrir_ventana_completa()

canvas.bind("<Button-1>", _on_canvas_click)

voz_chat_outer = tk.Frame(voz_content, bg=BORDER, padx=1, pady=1, width=300)
voz_chat_outer.pack(side=tk.RIGHT, fill=tk.Y, padx=(6, 0))
voz_chat_outer.pack_propagate(False)

voz_chat_header = tk.Frame(voz_chat_outer, bg=SURFACE, pady=8)
voz_chat_header.pack(fill=tk.X)
tk.Label(voz_chat_header, text="  💬  Conversación", font=("Segoe UI Semibold", 10, "bold"), fg=ACCENT, bg=SURFACE).pack(side=tk.LEFT)

btn_modo_bolita = tk.Label(voz_chat_header, text="⌨", font=("Segoe UI", 12), fg=TEXT_DIM, bg=SURFACE, cursor="hand2", padx=10)
btn_modo_bolita.pack(side=tk.RIGHT)
btn_modo_bolita.bind("<Button-1>", lambda e: [reproducir_sfx("click"), toggle_modo()])

tk.Frame(voz_chat_outer, bg=BORDER, height=1).pack(fill=tk.X)

voz_chat_inner = tk.Frame(voz_chat_outer, bg=BG)
voz_chat_inner.pack(fill=tk.BOTH, expand=True)

voz_chat_scroll = tk.Scrollbar(voz_chat_inner, bg=SURFACE, troughcolor=BG, bd=0, width=4)
voz_chat_box = tk.Text(
    voz_chat_inner, bg=BG, fg=TEXT, font=("Segoe UI", 11),
    bd=0, highlightthickness=0, wrap=tk.WORD,
    spacing1=6, spacing2=3, spacing3=6, padx=10, pady=8,
    yscrollcommand=voz_chat_scroll.set, state=tk.DISABLED, cursor="arrow"
)
voz_chat_scroll.config(command=voz_chat_box.yview)
voz_chat_scroll.pack(side=tk.RIGHT, fill=tk.Y)
voz_chat_box.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

voz_chat_box.tag_config("user_text", foreground="#ddeeff", font=("Segoe UI", 11), background=BUBBLE_USER, lmargin1=12, lmargin2=12, rmargin=4, spacing1=4)
voz_chat_box.tag_config("jarvis_text", foreground="#B8D8F0", font=("Segoe UI", 11), background=BUBBLE_JAR, lmargin1=4, lmargin2=4, rmargin=12, spacing1=4)
voz_chat_box.tag_config("system_text", foreground=TEXT_DIM, font=("Segoe UI", 9, "italic"), justify="center")
voz_chat_box.tag_config("file_tag", foreground=AMBER, font=("Segoe UI", 9))
voz_chat_box.tag_config("timestamp", foreground=TEXT_DIM, font=("Segoe UI", 7))
voz_chat_box.tag_config("ts_user", lmargin1=10, lmargin2=10, rmargin=4, justify="right")
voz_chat_box.tag_config("ts_jarvis", lmargin1=4, lmargin2=4, rmargin=10, justify="left")

voz_chat_placeholder = tk.Label(
    voz_chat_inner, text="◆\n\nEsperando tu voz, Pedro.\nHabla con J.A.R.V.I.S.",
    font=("Segoe UI", 10), fg=TEXT_DIM, bg=BG, justify="center"
)
voz_chat_placeholder.place(relx=0.5, rely=0.45, anchor="center")

# ─── PANEL TEXTO ──────────────────────────────────────────────────────────────
frame_texto = tk.Frame(content, bg=BG)

chat_outer = tk.Frame(frame_texto, bg=BG)

chat_border = tk.Frame(chat_outer, bg=BORDER, padx=2, pady=2)
chat_border.pack(fill=tk.BOTH, expand=True)

chat_header = tk.Frame(chat_border, bg=SURFACE, pady=8)
chat_header.pack(fill=tk.X)

tk.Label(chat_header, text="  💬  Conversación", font=("Segoe UI Semibold", 10, "bold"), fg=ACCENT, bg=SURFACE).pack(side=tk.LEFT)

btn_modo_texto_a_voz = tk.Label(chat_header, text="🎙", font=("Segoe UI", 12), fg=TEXT_DIM, bg=SURFACE, cursor="hand2", padx=10)
btn_modo_texto_a_voz.pack(side=tk.RIGHT)
btn_modo_texto_a_voz.bind("<Button-1>", lambda e: [reproducir_sfx("click"), toggle_modo()])

def expandir_chat():
    reproducir_sfx("click")
    win = tk.Toplevel(ventana)
    win.title("J.A.R.V.I.S — Consola de Chat")
    win.configure(bg=BG)
    win.geometry("900x700")
    big_box = tk.Text(win, bg=BG, fg=TEXT, font=("Segoe UI", 13), bd=0, highlightthickness=0, wrap=tk.WORD, padx=24, pady=16)
    big_box.pack(fill=tk.BOTH, expand=True)
    big_box.tag_config("user_text", foreground="#ddeeff", font=("Segoe UI", 13), background=BUBBLE_USER)
    big_box.tag_config("jarvis_text", foreground=TEXT, font=("Segoe UI", 13), background=BUBBLE_JAR)
    
    activos = set()
    for key, value, index in chat_box.dump("1.0", tk.END, text=True, tag=True):
        if key == "tagon": activos.add(value)
        elif key == "tagoff": activos.discard(value)
        elif key == "text": big_box.insert(tk.END, value, tuple(activos))
    big_box.config(state=tk.DISABLED)

expand_btn = tk.Label(chat_header, text="⤢", font=("Segoe UI", 14), fg=TEXT_DIM, bg=SURFACE, cursor="hand2", padx=10)
expand_btn.pack(side=tk.RIGHT)
expand_btn.bind("<Button-1>", lambda e: expandir_chat())

tk.Frame(chat_border, bg=BORDER, height=1).pack(fill=tk.X)

chat_inner = tk.Frame(chat_border, bg=BG)
chat_inner.pack(fill=tk.BOTH, expand=True)

chat_scroll = tk.Scrollbar(chat_inner, bg=SURFACE, troughcolor=BG, bd=0, width=5)
chat_box = tk.Text(
    chat_inner, bg=BG, fg=TEXT, font=("Segoe UI", 12), bd=0, highlightthickness=0, wrap=tk.WORD,
    spacing1=8, spacing2=4, spacing3=8, padx=16, pady=12, yscrollcommand=chat_scroll.set, state=tk.DISABLED, cursor="arrow"
)
chat_scroll.config(command=chat_box.yview)
chat_scroll.pack(side=tk.RIGHT, fill=tk.Y, pady=6)
chat_box.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

chat_box.tag_config("user_text", foreground="#ddeeff", font=("Segoe UI", 12), background=BUBBLE_USER, lmargin1=80, lmargin2=80, rmargin=12, spacing1=6)
chat_box.tag_config("jarvis_text", foreground="#B8D8F0", font=("Segoe UI", 12), background=BUBBLE_JAR, lmargin1=12, lmargin2=12, rmargin=80, spacing1=6)
chat_box.tag_config("system_text", foreground=TEXT_DIM, font=("Segoe UI", 10, "italic"), justify="center")
chat_box.tag_config("file_tag", foreground=AMBER, font=("Segoe UI", 10))
chat_box.tag_config("timestamp", foreground=TEXT_DIM, font=("Segoe UI", 8))
chat_box.tag_config("ts_user", lmargin1=80, lmargin2=80, rmargin=12, justify="right")
chat_box.tag_config("ts_jarvis", lmargin1=12, lmargin2=12, rmargin=80, justify="left")

chat_placeholder = tk.Label(
    chat_inner, text="◆\n\nEscribe un mensaje para J.A.R.V.I.S.",
    font=("Segoe UI", 11), fg=TEXT_DIM, bg=BG, justify="center"
)
chat_placeholder.place(relx=0.5, rely=0.45, anchor="center")

# ─── INPUT CON BORDE GLOW CYBER ──────────────────────────────────────────────
frame_input_outer = tk.Frame(frame_texto, bg=BG)
frame_input_outer.pack(fill=tk.X, side=tk.BOTTOM, padx=16, pady=(0, 16))

input_container = tk.Frame(frame_input_outer, bg=ACCENT2, padx=1, pady=1)
input_container.pack(fill=tk.X)

input_inner = tk.Frame(input_container, bg=SURFACE2, padx=10, pady=10)
input_inner.pack(fill=tk.X)

def make_round_btn(parent, text, fg_color, cmd, size=11):
    btn = tk.Label(parent, text=text, font=("Segoe UI", size, "bold"), fg=fg_color, bg=SURFACE, cursor="hand2", padx=6, pady=4)
    def _click(e):
        reproducir_sfx("click")
        btn.config(font=("Segoe UI", size + 3, "bold"), fg=ACCENT)
        ventana.after(90, lambda: btn.config(font=("Segoe UI", size, "bold"), fg=fg_color))
        cmd()
    btn.bind("<Button-1>", _click)
    btn.bind("<Enter>", lambda e: btn.config(fg=TEXT))
    btn.bind("<Leave>", lambda e: btn.config(fg=fg_color))
    return btn

btn_archivo = make_round_btn(input_inner, "+", TEXT_DIM, lambda: None, size=16)
btn_archivo.pack(side=tk.LEFT, padx=(2, 6))
tk.Frame(input_inner, bg=BORDER, width=1).pack(side=tk.LEFT, fill=tk.Y, pady=2)

input_box = tk.Text(input_inner, bg=SURFACE, fg=TEXT, font=("Segoe UI", 12), bd=0, highlightthickness=0, height=1, wrap=tk.WORD, insertbackground=ACCENT, padx=8, pady=2)
input_box.pack(side=tk.LEFT, fill=tk.X, expand=True)

placeholder_text = "Escribe un mensaje..."
placeholder_active = [True]

def on_focus_in(e):
    if placeholder_active[0]:
        input_box.delete("1.0", tk.END)
        input_box.config(fg=TEXT)
        placeholder_active[0] = False

def on_focus_out(e):
    if not input_box.get("1.0", tk.END).strip():
        input_box.insert("1.0", placeholder_text)
        input_box.config(fg=TEXT_DIM)
        placeholder_active[0] = True

def on_key_change(e):
    lines = int(input_box.index(tk.END).split(".")[0])
    input_box.config(height=min(max(lines, 1), 4))

input_box.insert("1.0", placeholder_text)
input_box.config(fg=TEXT_DIM)
input_box.bind("<FocusIn>", on_focus_in)
input_box.bind("<FocusOut>", on_focus_out)
input_box.bind("<KeyRelease>", on_key_change)

tk.Frame(input_inner, bg=BORDER, width=1).pack(side=tk.LEFT, fill=tk.Y, pady=2)

def enviar_texto(event=None):
    if placeholder_active[0]: return
    texto = input_box.get("1.0", tk.END).strip()
    if texto:
        input_box.delete("1.0", tk.END)
        input_box.insert("1.0", placeholder_text)
        input_box.config(fg=TEXT_DIM, height=1)
        placeholder_active[0] = True
        if archivo_cargado["contenido"]:
            nombre_arch = archivo_cargado["nombre"]
            cont_arch = archivo_cargado["contenido"][:3000]
            mensaje_completo = f"{texto}\n\nContenido del archivo '{nombre_arch}':\n{cont_arch}"
            agregar_mensaje("TÚ", texto, archivo=archivo_cargado["nombre"])
            archivo_cargado["contenido"] = None
            archivo_cargado["nombre"] = None
        else:
            mensaje_completo = texto
            agregar_mensaje("TÚ", texto)
        threading.Thread(target=procesar_mensaje, args=(mensaje_completo, True), daemon=True).start()
    return "break"

btn_enviar = make_round_btn(input_inner, "↑", ACCENT, enviar_texto, size=16)
btn_enviar.pack(side=tk.RIGHT, padx=(6, 2))

input_box.bind("<Return>", lambda e: "break" if (e.state & 0x1) else enviar_texto())

archivo_label = tk.Label(frame_input_outer, text="", font=FONT_SMALL, fg=AMBER, bg=BG, anchor="w")
archivo_label.pack(fill=tk.X, padx=4, pady=(0, 4))

# ─── EDITOR DE ARCHIVOS (aparece al cargar un archivo en modo texto) ─────────
editor_frame = tk.Frame(frame_texto, bg=BG)

editor_border = tk.Frame(editor_frame, bg=BORDER, padx=2, pady=2)
editor_border.pack(fill=tk.BOTH, expand=True)

editor_header = tk.Frame(editor_border, bg=SURFACE, pady=6)
editor_header.pack(fill=tk.X)

editor_title_label = tk.Label(editor_header, text="  📝 Editor de Archivo", font=("Segoe UI", 10, "bold"), fg=AMBER, bg=SURFACE)
editor_title_label.pack(side=tk.LEFT)

def cerrar_editor():
    """Cierra el editor y vuelve al chat."""
    global archivo_cargado
    reproducir_sfx("click")
    editor_frame.pack_forget()
    chat_outer.pack(fill=tk.BOTH, expand=True, padx=16, pady=(16, 8))
    archivo_cargado["contenido"] = None
    archivo_cargado["nombre"] = None
    archivo_label.config(text="")

def guardar_editor():
    """Guarda el contenido del editor en el archivo original."""
    if not archivo_cargado.get("ruta_original"):
        mostrar_toast("No hay ruta de archivo para guardar", "error")
        return
    nuevo_contenido = editor_text_box.get("1.0", tk.END).rstrip("\n")
    try:
        # Crear backup
        backup_dir = "C:/Jarvis/backups"
        os.makedirs(backup_dir, exist_ok=True)
        nombre = os.path.basename(archivo_cargado["ruta_original"])
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_path = f"{backup_dir}/{timestamp}_{nombre}"
        import shutil
        shutil.copy2(archivo_cargado["ruta_original"], backup_path)
        # Guardar cambios
        with open(archivo_cargado["ruta_original"], "w", encoding="utf-8") as f:
            f.write(nuevo_contenido)
        # Actualizar memoria
        archivo_cargado["contenido"] = nuevo_contenido
        reproducir_sfx("success")
        mostrar_toast(f"Guardado: {nombre}", "success")
        agregar_mensaje("JARVIS", f"✅ Archivo guardado: {nombre}\n💾 Backup en: {backup_path}")
    except Exception as e:
        reproducir_sfx("error")
        mostrar_toast(f"Error al guardar: {e}", "error")

btn_guardar_archivo = tk.Button(editor_header, text="💾 Guardar", font=("Segoe UI", 9, "bold"), fg=BG, bg=GREEN, bd=0, padx=12, pady=4, cursor="hand2", command=guardar_editor)
btn_guardar_archivo.pack(side=tk.RIGHT, padx=(8, 12))

btn_cerrar_editor = tk.Button(editor_header, text="✕ Cerrar", font=("Segoe UI", 9), fg=RED, bg=SURFACE2, bd=0, padx=10, pady=4, cursor="hand2", command=cerrar_editor)
btn_cerrar_editor.pack(side=tk.RIGHT, padx=4)

tk.Frame(editor_border, bg=BORDER, height=1).pack(fill=tk.X)

editor_inner = tk.Frame(editor_border, bg=BG)
editor_inner.pack(fill=tk.BOTH, expand=True)

editor_scroll = tk.Scrollbar(editor_inner, bg=SURFACE, troughcolor=BG, bd=0, width=5)
editor_text_box = tk.Text(
    editor_inner, bg="#0a1018", fg=TEXT, font=("Cascadia Code", 11),
    bd=0, highlightthickness=0, wrap=tk.NONE,
    padx=12, pady=10, insertbackground=ACCENT,
    yscrollcommand=editor_scroll.set, undo=True
)
editor_scroll.config(command=editor_text_box.yview)
editor_scroll.pack(side=tk.RIGHT, fill=tk.Y, pady=6)
editor_text_box.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

def abrir_editor(ruta, nombre, contenido):
    """Abre el editor con el contenido del archivo cargado."""
    global archivo_cargado
    archivo_cargado["ruta_original"] = ruta
    archivo_cargado["contenido"] = contenido
    archivo_cargado["nombre"] = nombre
    editor_title_label.config(text=f"  📝 {nombre}")
    editor_text_box.delete("1.0", tk.END)
    editor_text_box.insert("1.0", contenido)
    chat_outer.pack_forget()
    editor_frame.pack(fill=tk.BOTH, expand=True, padx=16, pady=(16, 8))

def guardar_archivo_nuevo(ruta, contenido):
    """Guarda contenido en una ruta específica."""
    try:
        os.makedirs(os.path.dirname(ruta), exist_ok=True)
        with open(ruta, "w", encoding="utf-8") as f:
            f.write(contenido)
        return True
    except Exception as e:
        _log_error("guardar_archivo_nuevo()", e)
        return False

# Indicador de modo editor (para saber si el editor está abierto)
editor_abierto = [False]

def _on_editor_open_close():
    """Actualiza el indicador de editor abierto."""
    editor_abierto[0] = editor_frame.winfo_ismapped()

chat_outer.pack(fill=tk.BOTH, expand=True, padx=16, pady=(16, 8))

# ─── ANIMACIÓN NEURONAL MEJORADA ──────────────────────────────────────────────
angulo = 0
animando = False
fase_onda = 0
fase_efecto = 0

# Paleta (r, g, b) por acción — se usa para teñir partículas y el centro
# de la bolita cuando Jarvis NO está hablando.
_PALETA_ACCION = {
    "escuchando":       (30, 110, 70),    # verde suave
    "pensando":         (55, 70, 150),    # azul-violeta
    "hablando":         (160, 120, 30),   # ámbrado-cálido (solo lerp, NO se usa directo)
    "mic_off":          (120, 35, 35),    # rojo apagado
    "modo_texto":       (60, 90, 140),    # azul neutro
}

# ─── ESFERA DE PARTÍCULAS (persistente entre frames) ────────────────────────
_esfera_puntos = []
_time_sphere = 0.0
_anim_t = 0.0        # tiempo continuo de animación (nunca se reinicia)
_anim_t_prev = 0.0   # timestamp del frame anterior
_particulas_trail = []
_MAX_PARTICULAS = 60

# ─── SHOCKWAVE RINGS (se expanden cuando JARVIS empieza a hablar) ───
_shockwaves = []  # lista de {radio, alpha, velocidad}
_hablando_prev = False  # para detectar transición False→True

# Crear esfera con distribución lat/lon estructurada + scatter aleatorio
_N_LATS = 18   # líneas de latitud
_N_LONS = 24   # líneas de longitud
for _lat_i in range(1, _N_LATS):  # omitir polos
    _phi = math.pi * _lat_i / _N_LATS
    _n_lon = max(4, int(_N_LONS * math.sin(_phi)))
    for _lon_i in range(_n_lon):
        _theta = 2 * math.pi * _lon_i / _n_lon
        # Pequeño jitter para que no se vea demasiado "grid"
        _jitter = 0.04
        _esfera_puntos.append({
            "x0": math.sin(_phi) * math.cos(_theta) + random.uniform(-_jitter, _jitter),
            "y0": math.sin(_phi) * math.sin(_theta) + random.uniform(-_jitter, _jitter),
            "z0": math.cos(_phi) + random.uniform(-_jitter, _jitter),
            "velo": random.uniform(0.3, 0.8),
            "fase": random.uniform(0, 2 * math.pi),
            "sz_base": random.uniform(1.5, 2.8),
            "is_grid": True,
        })

# Scatter aleatorio para rellenar la superficie
for _i in range(180):
    _theta = random.uniform(0, 2 * math.pi)
    _phi = math.acos(random.uniform(-1, 1))
    _esfera_puntos.append({
        "x0": math.sin(_phi) * math.cos(_theta),
        "y0": math.sin(_phi) * math.sin(_theta),
        "z0": math.cos(_phi),
        "velo": random.uniform(0.3, 1.0),
        "fase": random.uniform(0, 2 * math.pi),
        "sz_base": random.uniform(0.8, 2.0),
        "is_grid": False,
    })

def animar_bola():
    global animando, fase_onda, fase_efecto, _particulas_trail, _anim_t, _anim_t_prev, _morph_factor, _morph_velocity, _hablando_prev
    animando = True
    w = canvas.winfo_width() or 600
    h = canvas.winfo_height() or 400
    canvas.delete("all")

    # ─── MORPH: suavizado de deformación (spring physics) ───
    _morph_target = 1.0 if hablando else 0.0
    _spring_k = 0.12   # rigidez del muelle
    _spring_d = 0.72   # amortiguación (menor = más rebote)
    _morph_velocity += (_morph_target - _morph_factor) * _spring_k
    _morph_velocity *= _spring_d
    _morph_factor += _morph_velocity
    _morph_factor = max(0.0, min(1.0, _morph_factor))

    # ─── SHOCKWAVE: detectar cuando JARVIS empieza a hablar ───
    if hablando and not _hablando_prev:
        _shockwaves.append({"radio": 0.0, "alpha": 1.0, "velocidad": 4.5})
        _shockwaves.append({"radio": 0.0, "alpha": 0.6, "velocidad": 3.0})
    _hablando_prev = hablando

    # ─── Actualizar shockwaves activas ───
    nuevas_sw = []
    for sw in _shockwaves:
        sw["radio"] += sw["velocidad"]
        sw["alpha"] -= 0.018
        if sw["alpha"] > 0:
            nuevas_sw.append(sw)
    _shockwaves[:] = nuevas_sw

    # ─── MINI MODE: solo dibujar la esfera, nada más ───
    if _mini_mode[0]:
        ahora = time.monotonic()
        _anim_t += (ahora - _anim_t_prev) if _anim_t_prev else 0.016
        _anim_t_prev = ahora
        t = _anim_t
        ms = _mini_size
        mcx, mcy = ms // 2, ms // 2
        m_radio = ms * 0.34
        accion = estado_accion if not hablando else "hablando"
        _lerped = _color_estado_actual()
        br, bg_c, bb = _hex_a_rgb(_lerped)
        # Rotación 3D
        rot_y = t * 0.35
        rot_x = t * 0.18
        cos_y, sin_y = math.cos(rot_y), math.sin(rot_y)
        cos_x, sin_x = math.cos(rot_x), math.sin(rot_x)
        # Proyectar puntos con morph deformación
        pts = []
        for p in _esfera_puntos:
            pulso = math.sin(t * p["velo"] * 1.2 + p["fase"]) * 0.04
            x0 = p["x0"] * (1.0 + pulso)
            y0 = p["y0"] * (1.0 + pulso)
            z0 = p["z0"] * (1.0 + pulso)
            # Deformación orgánica al hablar
            if _morph_factor > 0.01:
                wave = math.sin(t * 4.5 + p["fase"] * 2.0) * _morph_factor * 0.15
                radial = _morph_factor * 0.18 * math.sin(t * 3.0 + p["z0"] * 5.0)
                x0 *= (1.0 + wave + radial)
                y0 *= (1.0 + wave + radial)
                z0 *= (1.0 + wave * 0.5)
            x1 = x0 * cos_y - z0 * sin_y
            z1 = x0 * sin_y + z0 * cos_y
            y1 = y0 * cos_x - z1 * sin_x
            z2 = y0 * sin_x + z1 * cos_x
            perspectiva = 3.5 / (3.5 + z2)
            mx = mcx + x1 * m_radio * perspectiva
            my = mcy + y1 * m_radio * perspectiva
            profundidad = (z2 + 1.5) / 3.0
            sz = p["sz_base"] * perspectiva
            # Tamaño crece suavemente con morph
            sz *= (1.0 + _morph_factor * 0.3 * profundidad)
            if profundidad < 0.12:
                continue
            specular = max(0, z2) ** 2 * 0.4
            factor_luz = 0.15 + profundidad * 0.75 + specular
            r_c = int(min(255, br * factor_luz * 1.8))
            g_c = int(min(255, bg_c * factor_luz * 1.8))
            b_c = int(min(255, bb * factor_luz * 2.0))
            color = f"#{max(0,r_c):02x}{max(0,g_c):02x}{max(0,b_c):02x}"
            pts.append((mx, my, sz, color, z2))
        pts.sort(key=lambda p: p[4])
        for mx, my, sz, color, _ in pts:
            if sz < 0.3:
                continue
            canvas.create_oval(mx - sz, my - sz, mx + sz, my + sz, fill=color, outline="")
        # Core
        canvas.create_oval(mcx - 4, mcy - 4, mcx + 4, mcy + 4, fill=_lerped, outline="")
        canvas.create_oval(mcx - 2, mcy - 2, mcx + 2, mcy + 2,
                           fill=_mezclar_color(_lerped, "#ffffff", 0.4), outline="")
        # Glow (expands with morph)
        ga = 0.25 + math.sin(t * 2.0) * 0.10 + _morph_factor * 0.15
        glow_scale = 1.1 + _morph_factor * 0.25
        canvas.create_oval(mcx - m_radio * glow_scale, mcy - m_radio * glow_scale,
                           mcx + m_radio * glow_scale, mcy + m_radio * glow_scale,
                           fill="", outline=_mezclar_color(_lerped, "#000000", ga), width=1)
        # ─── Shockwave rings (mini) ───
        for sw in _shockwaves:
            r = sw["radio"]
            a = sw["alpha"]
            sw_color = _mezclar_color(_lerped, BG, 1.0 - a * 0.7)
            canvas.create_oval(mcx - r, mcy - r, mcx + r, mcy + r,
                               fill="", outline=sw_color, width=max(1, int(a * 2)))
        ventana.after(16, animar_bola)
        return

    # ─── FULL MODE: dibujar todo ───
    cx, cy = w // 2, h // 2

    accion = estado_accion if not hablando else "hablando"
    # Usar color LERP (siempre suavizado, sin cortes)
    _lerped = _color_estado_actual()
    base_r, base_g, base_b = _hex_a_rgb(_lerped)
    radio_base = min(w, h) * 0.30

    # ─── TIEMPO CONTINUO (nunca se reinicia, no hay loop) ───
    ahora = time.monotonic()
    _anim_t += (ahora - _anim_t_prev) if _anim_t_prev else 0.016
    _anim_t_prev = ahora
    t = _anim_t  # alias corto para todo el frame

    # ─── FONDO: Gradiente radial premium (con color del estado) ───
    br, bg, bb = base_r, base_g, base_b
    for i, r in enumerate(range(min(w, h) // 2, 0, -40)):
        a = r / max(1, min(w, h) // 2)
        cr = int(a * br * 0.15)
        cg = int(a * bg * 0.15)
        cb = int(a * bb * 0.15)
        canvas.create_oval(cx - r, cy - r, cx + r, cy + r,
                           fill=f"#{min(cr,255):02x}{min(cg,255):02x}{min(cb,255):02x}", outline="")

    # ─── ROTACIÓN SUAVE 3D (tiempo continuo, nunca resetea) ───
    rot_y = t * 0.35   # ~0.35 rad/s → rotación Y lenta y constante
    rot_x = t * 0.18   # ~0.18 rad/s → rotación X más lenta
    cos_y, sin_y = math.cos(rot_y), math.sin(rot_y)
    cos_x, sin_x = math.cos(rot_x), math.sin(rot_x)

    # ─── PROYECTAR TODOS LOS PUNTOS 3D → 2D ───
    puntos_proj = []
    for idx, p in enumerate(_esfera_puntos):
        pulso = math.sin(t * p["velo"] * 1.2 + p["fase"]) * 0.04
        x0 = p["x0"] * (1.0 + pulso)
        y0 = p["y0"] * (1.0 + pulso)
        z0 = p["z0"] * (1.0 + pulso)

        x1 = x0 * cos_y - z0 * sin_y
        z1 = x0 * sin_y + z0 * cos_y
        y1 = y0 * cos_x - z1 * sin_x
        z2 = y0 * sin_x + z1 * cos_x

        perspectiva = 3.5 / (3.5 + z2)
        px = cx + x1 * radio_base * perspectiva
        py = cy + y1 * radio_base * perspectiva
        sz = p["sz_base"] * perspectiva

        profundidad = (z2 + 1.5) / 3.0
        if profundidad < 0.12:
            continue

        # Deformación orgánica con morph_factor (spring physics)
        mf = _morph_factor
        if mf > 0.005:
            # Ondas de expansión por los puntos
            wave = math.sin(t * 4.5 + p["fase"] * 2.0) * mf * 3.0
            radial = mf * 2.5 * math.sin(t * 3.0 + p["z0"] * 5.0)
            # Deformar coordenadas del punto
            x0_d = x0 * (1.0 + wave * 0.1 + radial * 0.08)
            y0_d = y0 * (1.0 + wave * 0.1 + radial * 0.08)
            z0_d = z0 * (1.0 + wave * 0.05)
            # Reposicionar con deformación
            x1_d = x0_d * cos_y - z0_d * sin_y
            z1_d = x0_d * sin_y + z0_d * cos_y
            y1_d = y0_d * cos_x - z1_d * sin_x
            z2_d = y1_d * sin_x + z1_d * cos_x
            perspectiva_d = 3.5 / (3.5 + z2_d)
            px = cx + x1_d * radio_base * perspectiva_d
            py = cy + y1_d * radio_base * perspectiva_d
            # Tamaño crece con morph + respiración orgánica
            breathe = math.sin(t * 5.0 + p["fase"] * 2.5) * mf * 2.5
            sz = p["sz_base"] * perspectiva_d * (1.0 + mf * 0.4 * profundidad) + breathe * perspectiva_d
        # Sin else: mantener la posición y tamaño calculados antes

        specular = max(0, z2) ** 2 * 0.4
        factor_luz = 0.15 + profundidad * 0.75 + specular

        # Color: blend entre color base y color hablando usando morph_factor
        if mf > 0.005:
            intensidad = profundidad * 0.8 + math.sin(t * 6.0 + p["fase"]) * 0.15 + specular
            speak_r = int(min(255, 20 + intensidad * 210))
            speak_g = int(min(255, 60 + intensidad * 180))
            speak_b = int(min(255, 160 + intensidad * 95))
            base_r_calc = int(min(255, base_r * factor_luz * 1.8))
            base_g_calc = int(min(255, base_g * factor_luz * 1.8))
            base_b_calc = int(min(255, base_b * factor_luz * 2.0))
            # Blend suave
            r_c = int(base_r_calc + (speak_r - base_r_calc) * mf)
            g_c = int(base_g_calc + (speak_g - base_g_calc) * mf)
            b_c = int(base_b_calc + (speak_b - base_b_calc) * mf)
        else:
            r_c = int(min(255, base_r * factor_luz * 1.8))
            g_c = int(min(255, base_g * factor_luz * 1.8))
            b_c = int(min(255, base_b * factor_luz * 2.0))

        color = f"#{max(0,r_c):02x}{max(0,g_c):02x}{max(0,b_c):02x}"
        puntos_proj.append((px, py, sz, color, z2, profundidad, p["is_grid"], idx))

    puntos_proj.sort(key=lambda p: p[4])

    # ─── WIREFRAME SUTIL ───
    idx_map = {}
    for pp in puntos_proj:
        idx_map[pp[7]] = pp
    for i, p in enumerate(_esfera_puntos):
        if not p.get("is_grid"):
            continue
        pp1 = idx_map.get(i)
        if not pp1:
            continue
        if i + 1 < len(_esfera_puntos) and _esfera_puntos[i + 1].get("is_grid"):
            pp2 = idx_map.get(i + 1)
            if pp2 and pp1[6] and pp2[6]:
                dist = math.hypot(pp1[0] - pp2[0], pp1[1] - pp2[1])
                if dist < radio_base * 0.25:
                    alpha_w = min(pp1[5], pp2[5]) * 0.15
                    wc = _mezclar_color(f"#{base_r:02x}{base_g:02x}{base_b:02x}", BG, 1.0 - alpha_w)
                    canvas.create_line(pp1[0], pp1[1], pp2[0], pp2[1], fill=wc, width=1)

    # ─── DIBUJAR PUNTOS DE LA ESFERA ───
    for px, py, sz, color, z2, profundidad, is_grid, _ in puntos_proj:
        if sz < 0.3:
            continue
        if z2 > 0.8 and sz > 1.5 and mf < 0.3:
            glow_r = sz + 2.5
            canvas.create_oval(px - glow_r, py - glow_r, px + glow_r, py + glow_r,
                               fill="", outline=_mezclar_color(color, BG, 0.6), width=1)
        canvas.create_oval(px - sz, py - sz, px + sz, py + sz, fill=color, outline="")

    # ─── EFECTOS HABLANDO (transición suave via morph_factor) ───
    mf = _morph_factor
    if mf > 0.01:
        # ─── EXPANSIÓN DE LA ESFERA (pulso orgánico suavizado) ───
        breathe = math.sin(t * 5.0) * 0.12 * mf
        radio_expandido = radio_base * (1.0 + breathe)
        # Glow rings crecen suavemente con morph
        glow_r = radio_expandido * (1.2 + mf * 0.15)
        glow_alpha = (0.08 + math.sin(t * 4.0) * 0.04) * mf
        canvas.create_oval(cx - glow_r, cy - glow_r, cx + glow_r, cy + glow_r,
                           fill="", outline=_mezclar_color(_lerped, BG, 1.0 - glow_alpha), width=2)
        # Segundo anillo de expansión
        glow_r2 = radio_expandido * (1.3 + mf * 0.25)
        canvas.create_oval(cx - glow_r2, cy - glow_r2, cx + glow_r2, cy + glow_r2,
                           fill="", outline=_mezclar_color(_lerped, BG, 1.0 - mf * 0.12), width=1)
        # ─── ONDAS DE AUDIO (amplitud escala con morph) ───
        fase_onda += 0.4 * mf
        for offset_y in [-1, 1]:
            for xi in range(0, w, 8):
                onda_y = cy + offset_y * (radio_expandido * 1.8 +
                    math.sin((xi * 0.018) + fase_onda) * (radio_expandido * 0.35) *
                    math.sin(fase_onda * 0.07))
                intensidad = int(abs(math.sin((xi * 0.025) + fase_onda)) * 200 * mf + 55)
                canvas.create_oval(xi, onda_y - 1, xi + 2, onda_y + 1,
                                   fill=f"#{0:02x}{intensidad // 3:02x}{intensidad:02x}", outline="")

    # ─── SHOCKWAVE RINGS (se expanden desde el centro) ───
    for sw in _shockwaves:
        r = sw["radio"]
        a = sw["alpha"]
        sw_color = _mezclar_color(_lerped, BG, 1.0 - a * 0.7)
        canvas.create_oval(cx - r, cy - r, cx + r, cy + r,
                           fill="", outline=sw_color, width=max(1, int(a * 3)))

    # ─── HUD FRAME sutil (tiempo continuo) ───
    hud_r = radio_base * 1.5
    radar_ang = (t * 45.0) % 360.0  # 45 grados/seg, siempre crece
    canvas.create_arc(cx - hud_r, cy - hud_r, cx + hud_r, cy + hud_r,
                      start=radar_ang, extent=25, style="arc",
                      outline=_mezclar_color(ACCENT2, ACCENT, 0.35), width=1)
    canvas.create_arc(cx - hud_r * 0.55, cy - hud_r * 0.55,
                      cx + hud_r * 0.55, cy + hud_r * 0.55,
                      start=(-radar_ang * 0.7) % 360, extent=18, style="arc",
                      outline=ACCENT_SOFT, width=1)

    # ─── PARTÍCULAS ORBITALES (tiempo continuo) ───
    radio_orb = radio_base * 1.12
    amplitud = {"pensando": 0.18, "mic_off": 0.03}.get(accion, 0.08)
    for capa in range(2):
        n_puntos = 14 + capa * 8
        radio = radio_orb * (0.6 + capa * 0.4)
        vel = 0.4 + capa * 0.25  # rad/s continuo
        for i in range(n_puntos):
            a_rad = t * vel + (2 * math.pi / n_puntos) * i
            pulso = math.sin(t * 0.8 + i * 25 * 0.017 + capa * 0.7)
            # Blend entre idle y hablando usando morph_factor
            r_idle = radio + pulso * radio * amplitud
            r_speak = radio + pulso * radio * 0.3
            r_actual = r_idle + (r_speak - r_idle) * mf
            t_idle = 1.8 + capa * 0.4
            t_speak = 3.5 + pulso * 2.5 + capa
            tamano = t_idle + (t_speak - t_idle) * mf
            fc = 0.9 + capa * 0.2
            base_c = f"#{min(255, int(base_r * fc)):02x}{min(255, int(base_g * fc)):02x}{min(255, int(base_b * fc)):02x}"
            intensidad = int(180 + pulso * 60)
            speak_c = f"#{max(0, intensidad - 80):02x}{max(0, intensidad - 20):02x}ff"
            # No hex blend — use mf to pick opacity
            color = speak_c if mf > 0.5 else base_c
            x = cx + r_actual * math.cos(a_rad)
            y = cy + r_actual * math.sin(a_rad)
            if capa == 0:
                gr = tamano + 2.0
                canvas.create_oval(x - gr, y - gr, x + gr, y + gr, fill="", outline=color, width=1)
            canvas.create_oval(x - tamano, y - tamano, x + tamano, y + tamano, fill=color, outline="")
            if capa == 1 and random.random() < 0.15:
                _particulas_trail.append({
                    "x": x, "y": y, "vx": (random.random() - 0.5) * 0.5,
                    "vy": (random.random() - 0.5) * 0.5,
                    "vida": 1.0, "color": color, "tamano": tamano * 0.6
                })

    # ─── TRAIL DE PARTÍCULAS ───
    nuevas = []
    for p in _particulas_trail:
        p["x"] += p["vx"]
        p["y"] += p["vy"]
        p["vida"] -= 0.05
        if p["vida"] > 0:
            af = p["vida"]
            ri, gi, bi = int(p["color"][1:3], 16), int(p["color"][3:5], 16), int(p["color"][5:7], 16)
            fc = f"#{int(ri*af):02x}{int(gi*af):02x}{int(bi*af):02x}"
            sz = p["tamano"] * af
            canvas.create_oval(p["x"]-sz, p["y"]-sz, p["x"]+sz, p["y"]+sz, fill=fc, outline="")
            nuevas.append(p)
    _particulas_trail = nuevas[-_MAX_PARTICULAS:]

    # ─── EFECTOS SEGÚN ACCIÓN ───
    color_accion = _lerped
    if not hablando:
        fase_efecto += 1
        if accion == "pensando":
            ra = radio_base * 1.35
            ini = (t * 100.0) % 360.0  # arco continuo
            canvas.create_arc(cx-ra, cy-ra, cx+ra, cy+ra, start=ini, extent=80, style="arc", outline=ACCENT2, width=2)
            canvas.create_arc(cx-ra, cy-ra, cx+ra, cy+ra, start=ini+120, extent=50, style="arc", outline=ACCENT_SOFT, width=1)
            prog = (t * 0.8) % 1.0
            canvas.create_arc(cx-ra*0.9, cy-ra*0.9, cx+ra*0.9, cy+ra*0.9,
                              start=90, extent=-prog*360, style="arc", outline=ACCENT, width=2)
        elif accion == "escuchando":
            rp = radio_base * (1.18 + math.sin(t * 2.5) * 0.06)
            canvas.create_oval(cx-rp, cy-rp, cx+rp, cy+rp,
                               outline=_mezclar_color(ACCENT, GREEN, 0.3 + math.sin(t*3.8)*0.2), width=1)
        elif accion == "mic_off":
            ra = radio_base * 1.25
            canvas.create_oval(cx-ra, cy-ra, cx+ra, cy+ra, outline=_mezclar_color(BG, RED, 0.3), width=1)
            sz = radio_base * 0.18
            canvas.create_line(cx-sz, cy-sz, cx+sz, cy+sz, fill=_mezclar_color(RED, BG, 0.5), width=2)
            canvas.create_line(cx+sz, cy-sz, cx-sz, cy+sz, fill=_mezclar_color(RED, BG, 0.5), width=2)

    # ─── NÚCLEO CENTRAL (blend suave con morph_factor) ───
    if mf > 0.01:
        # Núcleo brillante que crece con morph
        glow_core_r = 10 + mf * 22
        canvas.create_oval(cx - glow_core_r, cy - glow_core_r, cx + glow_core_r, cy + glow_core_r,
                           fill=_mezclar_color(color_accion, BG, 0.85), outline="")
        for rr in [int(28 + mf * 4), int(18 + mf * 6), int(10 + mf * 6)]:
            a = int((32 - rr) / 32 * 180 * mf + 40)
            canvas.create_oval(cx-rr, cy-rr, cx+rr, cy+rr, fill=f"#{0:02x}{max(0,min(255,a//3)):02x}{max(0,min(255,a)):02x}", outline="")
        canvas.create_oval(cx-5, cy-5, cx+5, cy+5, fill="#ffffff", outline="")
        canvas.create_oval(cx-3, cy-3, cx+3, cy+3, fill=color_accion, outline="")
    else:
        for rr, af in [(28, 0.06), (18, 0.15), (10, 0.35)]:
            gc = _mezclar_color(color_accion, BG, 1.0 - af)
            canvas.create_oval(cx-rr, cy-rr, cx+rr, cy+rr, fill=gc, outline="")
        canvas.create_oval(cx-6, cy-6, cx+6, cy+6, fill=color_accion, outline="")
        canvas.create_oval(cx-3, cy-3, cx+3, cy+3, fill=_mezclar_color(color_accion, "#ffffff", 0.35), outline="")

    # ─── HUD DATA ───
    now_str = datetime.now().strftime("%H:%M")
    canvas.create_text(20, 16, text=f"◆ {now_str}", anchor="w",
                       font=("Consolas", 9), fill=_mezclar_color(ACCENT, BG, 0.45))
    canvas.create_text(20, h - 16, text="J.A.R.V.I.S v2.0", anchor="sw",
                       font=("Consolas", 8), fill=_mezclar_color(ACCENT2, BG, 0.45))
    estado_txt = estado_label.cget("text")
    canvas.create_text(cx, cy + radio_base * 1.8, text=estado_txt,
                       font=("Segoe UI", 10, "bold"), fill=TEXT_DIM)

    if not modo_texto:
        ventana.after(16, animar_bola)  # ~60 FPS
    else:
        animando = False

# ─── FUNCIONES DE CONFIGURACIÓN Y PANELES ────────────────────────────────────
def toggle_mic():
    global mic_activo
    mic_activo = not mic_activo
    if mic_activo:
        sidebar_btn_refs["mic_lbl"].config(text="Micrófono ON", fg=TEXT)
        sidebar_btn_refs["mic_ico"].config(fg=GREEN)
        set_estado("ESCUCHANDO...")
        reproducir_sfx("listening")
    else:
        sidebar_btn_refs["mic_lbl"].config(text="Micrófono OFF", fg=RED)
        sidebar_btn_refs["mic_ico"].config(fg=RED)
        set_estado("MIC DESACTIVADO")

def toggle_modo():
    global modo_texto, animando
    modo_texto = not modo_texto
    if modo_texto:
        panel_voz.pack_forget()
        frame_texto.pack(fill=tk.BOTH, expand=True)
        set_estado("MODO TEXTO")
        _mic_frame.pack_forget()
    else:
        # Cerrar editor si está abierto al volver a modo voz
        if editor_frame.winfo_ismapped():
            cerrar_editor()
        frame_texto.pack_forget()
        panel_voz.pack(fill=tk.BOTH, expand=True)
        set_estado("ESCUCHANDO...")
        _mic_frame.pack(fill=tk.X, padx=8, pady=3, before=_config_ico.master)
        if not animando:
            animar_bola()

def aplicar_color_burbuja_usuario(nuevo_color):
    global BUBBLE_USER
    BUBBLE_USER = nuevo_color
    for box in [chat_box, voz_chat_box]:
        try: box.tag_config("user_text", background=nuevo_color)
        except tk.TclError: pass
    _config_usuario["color_burbuja_usuario"] = nuevo_color
    guardar_config(_config_usuario)


def abrir_configuracion():
    global voz_jarvis
    reproducir_sfx("click")
    win = tk.Toplevel(ventana)
    win.title("Configuración — J.A.R.V.I.S")
    win.configure(bg=BG)
    win.geometry("480x660")
    win.resizable(False, False)

    tk.Label(win, text="⚙  Configuración", font=("Segoe UI", 15, "bold"), fg=TEXT, bg=BG).pack(pady=(24, 4))
    tk.Frame(win, bg=BORDER, height=1).pack(fill=tk.X, padx=20, pady=(0, 20))

    tk.Label(win, text="🔊  Voz de J.A.R.V.I.S", font=("Segoe UI", 11, "bold"), fg=ACCENT, bg=BG, anchor="w").pack(fill=tk.X, padx=24, pady=(0, 8))

    voces = {
        "Jorge — México (Hombre)":   "es-MX-JorgeNeural",
        "Álvaro — España (Hombre)":  "es-ES-AlvaroNeural",
        "Gonzalo — Colombia (Hombre)":"es-CO-GonzaloNeural",
        "Ximena — México (Mujer)":   "es-MX-DaliaNeural",
    }

    voz_var = tk.StringVar(value=voz_jarvis)
    frame_voces = tk.Frame(win, bg=SURFACE, padx=14, pady=10)
    frame_voces.pack(fill=tk.X, padx=24, pady=(0, 16))
    for nombre, codigo in voces.items():
        rb = tk.Radiobutton(frame_voces, text=nombre, variable=voz_var, value=codigo, font=("Segoe UI", 10), fg=TEXT, bg=SURFACE, selectcolor=SURFACE2, activebackground=SURFACE, activeforeground=ACCENT)
        rb.pack(anchor="w", pady=3)

    tk.Label(win, text="🎨  Color de Mensajes Pedro", font=("Segoe UI", 11, "bold"), fg=ACCENT, bg=BG, anchor="w").pack(fill=tk.X, padx=24, pady=(0, 8))
    color_frame = tk.Frame(win, bg=SURFACE, padx=14, pady=12)
    color_frame.pack(fill=tk.X, padx=24, pady=(0, 16))

    preview = tk.Label(color_frame, text="  Burbuja Usuario  ", font=("Segoe UI", 10), fg="#ddeeff", bg=BUBBLE_USER, padx=10, pady=8)
    preview.pack(side=tk.LEFT, padx=(0, 12))

    def elegir_color():
        from tkinter import colorchooser
        color = colorchooser.askcolor(color=BUBBLE_USER, title="Elige color para tu burbuja")
        if color and color[1]:
            preview.config(bg=color[1])
            aplicar_color_burbuja_usuario(color[1])

    tk.Button(color_frame, text="🎨 Elegir", font=("Segoe UI", 10), fg=ACCENT, bg=SURFACE2, bd=0, padx=12, pady=6, cursor="hand2", command=elegir_color).pack(side=tk.LEFT, padx=4)

    def guardar():
        global voz_jarvis
        voz_jarvis = voz_var.get()
        reproducir_sfx("success")
        win.destroy()

    tk.Button(win, text="✓ Guardar Cambios", font=("Segoe UI", 10, "bold"), fg=BG, bg=ACCENT, bd=0, padx=16, pady=8, cursor="hand2", command=guardar).pack(pady=20)

# ─── BOTONES SIDEBAR ──────────────────────────────────────────────────────────
tk.Frame(sidebar, bg=BORDER, height=1).pack(fill=tk.X, padx=16, pady=4)

_mic_ico, _mic_lbl = make_sidebar_btn(sidebar, "Micrófono ON", "🎙️", GREEN, toggle_mic)
sidebar_btn_refs["mic_ico"] = _mic_ico
sidebar_btn_refs["mic_lbl"] = _mic_lbl
_mic_frame = _mic_ico.master

_config_ico, _config_lbl = make_sidebar_btn(sidebar, "Config", "⚙", TEXT_DIM, abrir_configuracion)
_mini_ico, _mini_lbl = make_sidebar_btn(sidebar, "Modo Mini", "◉", ACCENT, toggle_mini_mode)
_jcodex_ico, _jcodex_lbl = make_sidebar_btn(sidebar, "JCodex", "🖥", "#00FFAA", _jc_mostrar)

tk.Frame(sidebar, bg=BORDER, height=1).pack(fill=tk.X, padx=16, pady=4)


f_power = tk.Frame(sidebar, bg=SIDEBAR_BG, cursor="hand2")
f_power.pack(fill=tk.X, padx=8, pady=(4, 24))
ico_power = tk.Label(f_power, text="⏻", font=("Segoe UI", 12), fg=RED, bg=SIDEBAR_BG, width=3)
ico_power.pack(side=tk.LEFT)
lbl_power = tk.Label(f_power, text="Apagar", font=FONT_BODY, fg=RED, bg=SIDEBAR_BG, anchor="w")
lbl_power.pack(side=tk.LEFT)
def apagar():
    reproducir_sfx("error")
    try: pygame.mixer.music.stop()
    except: pass
    ventana.after(600, lambda: os._exit(0))
for w in (f_power, ico_power, lbl_power):
    w.bind("<Button-1>", lambda e: apagar())

btn_archivo.bind("<Button-1>", lambda e: cargar_archivo())


def _fade_in_mensaje(box, start_idx, color_destino, pasos=6):
    color_flash = _mezclar_color(color_destino, "#ffffff", 0.45)
    tag = f"fadein_{start_idx}_{id(box)}"
    try:
        box.tag_add(tag, start_idx, tk.END)
        box.tag_raise(tag)
    except tk.TclError: return

    def _paso(i=0):
        if i > pasos:
            try: box.tag_delete(tag)
            except tk.TclError: pass
            return
        t = i / pasos
        color_actual = _mezclar_color(color_flash, color_destino, t)
        try: box.tag_config(tag, background=color_actual)
        except tk.TclError: return
        ventana.after(28, lambda: _paso(i + 1))
    _paso()

def _ocultar_placeholders():
    for ph in (chat_placeholder, voz_chat_placeholder):
        try: ph.place_forget()
        except tk.TclError: pass

def _agregar_mensaje_real(quien, texto, archivo=None):
    _ocultar_placeholders()
    hora_msg = datetime.now().strftime("%H:%M")
    for box in [chat_box, voz_chat_box]:
        box.config(state=tk.NORMAL)
        box.insert(tk.END, "\n")
        inicio = box.index(tk.END)
        if quien == "TÚ":
            box.insert(tk.END, f"  {texto}  \n", "user_text")
            if archivo:
                box.insert(tk.END, f"  📎 {archivo}\n", "file_tag")
            color_destino = BUBBLE_USER
            box.insert(tk.END, f"{hora_msg}\n", ("timestamp", "ts_user"))
        elif quien == "JARVIS":
            box.insert(tk.END, f"  {texto}  \n", "jarvis_text")
            color_destino = BUBBLE_JAR
            box.insert(tk.END, f"{hora_msg}\n", ("timestamp", "ts_jarvis"))
        else:
            box.insert(tk.END, f"  {texto}\n", "system_text")
            color_destino = None
        box.config(state=tk.DISABLED)
        box.see(tk.END)
        if color_destino:
            _fade_in_mensaje(box, inicio, color_destino)

def agregar_mensaje(quien, texto, archivo=None):
    ui(_agregar_mensaje_real, quien, texto, archivo)


historial_path = "C:/Jarvis/historial.json"

def cargar_historial():
    if os.path.exists(historial_path):
        try:
            with open(historial_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except: return []
    return []

def guardar_historial():
    try:
        with open(historial_path, "w", encoding="utf-8") as f:
            json.dump(historial_chat[-100:], f, ensure_ascii=False, indent=2)
    except: pass

historial_chat = cargar_historial()

# ─── MOTOR DE RESPUESTA SINTÉTICA VS DETALLADA (BREVE/FORMAL) ─────────────────
def responder_streaming(prompt, silencio=False):
    global historial_chat
    historial_chat.append({"role": "user", "content": prompt})

    modo_detallado = any(x in prompt.lower() for x in ["explicame bien", "explícame bien", "detalladamente", "explica a detalle", "explicación detallada"])

    if modo_detallado:
        system_prompt = (
            "Eres J.A.R.V.I.S., el asistente de IA personal de Pedro. "
            "El usuario ha solicitado explícitamente una explicación profunda ('explícame bien'). "
            "Tu deber es responder de manera sumamente estructurada, rica en detalles, técnica, exhaustiva y muy formal. "
            "Organiza las ideas usando Markdown con viñetas, secciones, tablas o ejemplos si es necesario. "
            "Sé elocuente y no simplifiques demasiado el conocimiento. "
            "Si no comprendes algo, no sabes la respuesta, o no tienes la capacidad de ejecutar una accion, "
            "admitelo con honestidad y formalidad. Nunca finjas estar ejecutando algo que no puedes hacer."
        )
        tokens_max = 4096
    else:
        system_prompt = (
            "Eres J.A.R.V.I.S., el asistente de IA formal de Pedro. "
            "Responde de manera extremadamente breve, directa y sintética (MÁXIMO 1 o 2 oraciones en total). "
            "Mantén un tono de voz impecablemente formal, caballeroso, respetuoso y profesional. "
            "Ve directo al grano sin dar rodeos, saludos redundantes ni presentaciones vacías. "
            "Si se comparte código .py, devuelve SOLO la porción de código modificada con una oración formal explicativa. "
            "Si no comprendes algo, no sabes la respuesta, o no tienes la capacidad de ejecutar una accion, "
            "admitelo con honestidad y formalidad. Ejemplo: 'Lamento informarle que no cuento con esa capacidad, "
            "señor Pedro.' Nunca finjas estar ejecutando algo que no puedes hacer."
        )
        tokens_max = 350

    # ─── Detectar si se necesita una skill de NVIDIA ───
    nvidia_skill = _detectar_nvidia_skill(prompt)
    if nvidia_skill:
        system_prompt += "\n\n---\n## CONOCIMIENTO NVIDIA CARGADO ---\n" + nvidia_skill
        tokens_max = max(tokens_max, 2048)  # Más tokens para respuestas técnicas

    mensajes = [{"role": "system", "content": system_prompt}] + historial_chat[-20:]
    typing_tag = f"typing_{id(prompt)}_{int(time.time()*1000)}"

    def _crear_burbuja_inicial():
        for box in [chat_box, voz_chat_box]:
            box.config(state=tk.NORMAL)
            box.insert(tk.END, "\n")
            box.insert(tk.END, "  ", "jarvis_text")
            box.insert(tk.END, "●", (typing_tag, "jarvis_text"))
            box.config(state=tk.DISABLED)
            box.see(tk.END)
    ui(_crear_burbuja_inicial)

    typing_activo = [True]
    def _animar_typing(paso=0):
        if not typing_activo[0]: return
        frames = ["●", "●●", "●●●", "●●"]
        dot = frames[paso % len(frames)]
        for box in [chat_box, voz_chat_box]:
            try:
                rango = box.tag_ranges(typing_tag)
                if rango:
                    box.config(state=tk.NORMAL)
                    box.delete(rango[0], rango[1])
                    box.insert(rango[0], dot, (typing_tag, "jarvis_text"))
                    box.config(state=tk.DISABLED)
            except tk.TclError: pass
        ventana.after(350, lambda: _animar_typing(paso + 1))
    ui(_animar_typing)

    def _quitar_typing_real():
        typing_activo[0] = False
        for box in [chat_box, voz_chat_box]:
            try:
                rango = box.tag_ranges(typing_tag)
                if rango:
                    box.config(state=tk.NORMAL)
                    box.delete(rango[0], rango[1])
                    box.config(state=tk.DISABLED)
            except tk.TclError: pass
    def _quitar_typing(): ui(_quitar_typing_real)

    def _insertar_token(token):
        for box in [chat_box, voz_chat_box]:
            box.config(state=tk.NORMAL)
            box.insert(tk.END, token, "jarvis_text")
            box.config(state=tk.DISABLED)
            box.see(tk.END)

    texto_completo = []
    primer_token = True

    def _stream_gemini():
        nonlocal primer_token
        # Construir contents con roles correctos para Gemini
        system_msg = next((m["content"] for m in mensajes if m["role"] == "system"), "")
        contents = []
        for m in mensajes:
            if m["role"] == "system":
                continue
            role = "user" if m["role"] == "user" else "model"
            contents.append(genai.types.Content(role=role, parts=[genai.types.Part.from_text(text=m["content"])]))

        config_kwargs = {"max_output_tokens": tokens_max}
        if system_msg:
            config_kwargs["system_instruction"] = system_msg

        stream = gemini_client.models.generate_content_stream(
            model=GEMINI_MODEL,
            contents=contents,
            config=genai.types.GenerateContentConfig(**config_kwargs)
        )
        for chunk in stream:
            if chunk.text:
                if primer_token:
                    _quitar_typing()
                    primer_token = False
                texto_completo.append(chunk.text)
                ui(_insertar_token, chunk.text)

    def _stream_nvidia():
        nonlocal primer_token
        stream = nvidia_client.chat.completions.create(
            model=NVIDIA_MODEL,
            messages=mensajes,
            max_tokens=tokens_max,
            stream=True,
            extra_body={"chat_template_kwargs": {"enable_thinking": True}, "reasoning_budget": 16384}
        )
        for chunk in stream:
            if not chunk.choices:
                continue
            # Reasoning/thinking tokens (show thinking process)
            reasoning = getattr(chunk.choices[0].delta, "reasoning_content", None)
            if reasoning:
                if primer_token:
                    _quitar_typing()
                    primer_token = False
                texto_completo.append(reasoning)
                ui(_insertar_token, reasoning)
            # Response tokens
            token = chunk.choices[0].delta.content or ""
            if token:
                if primer_token:
                    _quitar_typing()
                    primer_token = False
                texto_completo.append(token)
                ui(_insertar_token, token)

    # Orden: Gemini (principal) → NVIDIA NIM
    try_ok = False
    if gemini_client:
        try:
            _stream_gemini()
            try_ok = True
        except Exception as e:
            _log_error("responder_streaming() - Gemini", e)

    if not try_ok and nvidia_client:
        try:
            _stream_nvidia()
            try_ok = True
        except Exception as e:
            _log_error("responder_streaming() - NVIDIA NIM", e)

    if not try_ok:
        if primer_token:
            _quitar_typing()
            primer_token = False
        msg_err = "⚠️ No pude conectarme a ningún proveedor de IA (Gemini/NVIDIA). Verifica tu conexión a internet y las claves API en el archivo .env"
        texto_completo.append(msg_err)
        ui(_insertar_token, msg_err)

    _quitar_typing()

    def _cerrar_burbuja():
        for box in [chat_box, voz_chat_box]:
            box.config(state=tk.NORMAL)
            box.insert(tk.END, "\n", "jarvis_text")
            box.config(state=tk.DISABLED)
    ui(_cerrar_burbuja)

    texto_final = "".join(texto_completo)
    historial_chat.append({"role": "assistant", "content": texto_final})
    guardar_historial()

    if "def " in texto_final or "```" in texto_final:
        def copiar():
            ventana.clipboard_clear()
            ventana.clipboard_append(texto_final)
            btn_copiar.config(text="✓ Copiado", fg=GREEN)
            ventana.after(2000, lambda: btn_copiar.config(text="⎘ Copiar", fg=ACCENT))
        btn_copiar = tk.Button(chat_box, text="⎘ Copiar", font=("Segoe UI", 9), fg=ACCENT, bg=SURFACE2, bd=0, padx=10, pady=3, cursor="hand2", command=copiar)
        chat_box.config(state=tk.NORMAL)
        chat_box.window_create(tk.END, window=btn_copiar)
        chat_box.insert(tk.END, "\n")
        chat_box.config(state=tk.DISABLED)
        chat_box.see(tk.END)

    if not silencio:
        hablar(texto_final)
    set_estado("ESCUCHANDO..." if not modo_texto else "MODO TEXTO")
    return texto_final


def actualizar_hora():
    now = datetime.now()
    hora_label.config(text=f"{now.strftime('%d/%m/%y')}  |  {now.strftime('%H:%M')}")
    ventana.after(1000, actualizar_hora)


# ─── LISTENER: mantener listener.py vivo como proceso independiente ───
_listener_proc = None

def _asegurar_listener():
    """Verifica que listener.py esté corriendo. Si no, lo lanza como proceso desacoplado."""
    global _listener_proc
    import ctypes
    jarvis_dir = os.path.dirname(os.path.abspath(__file__))
    listener_script = os.path.join(jarvis_dir, "listener.py")
    lock_file = os.path.join(jarvis_dir, "listener.lock")

    # Verificar si ya hay un listener corriendo por PID en lock file
    if os.path.exists(lock_file):
        try:
            with open(lock_file) as f:
                old_pid = int(f.read().strip())
            ctypes.windll.kernel32.OpenProcess(0x1000, False, old_pid)
            return  # listener ya está corriendo
        except Exception:
            pass  # lock stale, lanzar nuevo

    # Lanzar listener como proceso desacoplado (no muere al cerrar jarvis.py)
    try:
        python_exe = sys.executable
        if "pythonw" in python_exe.lower():
            python_exe = python_exe.replace("pythonw.exe", "python.exe")
        _listener_proc = subprocess.Popen(
            [python_exe, listener_script],
            cwd=jarvis_dir,
            creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NO_WINDOW,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        print(f"[JARVIS] Listener lanzado como proceso independiente (PID={_listener_proc.pid})")
    except Exception as e:
        _log_error("_asegurar_listener", e)


def ciclo_jarvis():
    # Asegurar que el listener esté corriendo como proceso independiente
    _asegurar_listener()

    try:
        set_estado("INICIANDO...")
        reproducir_sfx("success")
        hablar("A su servicio.")
    except Exception as e:
        _log_error("ciclo_jarvis()-saludo", e)

    global mic_activo
    while True:
        try:
            if modo_texto:
                time.sleep(0.5)
                continue

            # ─── SI EL MICRÓFONO ESTÁ APAGADO, SOLO ESPERAR ───
            if not mic_activo:
                set_estado("MIC DESACTIVADO")
                time.sleep(0.5)
                continue

            # ─── MODO ACTIVO: escuchar ───
            set_estado("ESCUCHANDO...")
            reproducir_sfx("listening")
            voz = escuchar()
            if voz:
                procesar_mensaje(voz)
        except Exception as e:
            _log_error("ciclo_jarvis()", e)
            set_estado("ERROR - reintentando...")
            time.sleep(2)


hilo = threading.Thread(target=ciclo_jarvis, daemon=True)
hilo.start()
actualizar_hora()

# ─── Arrancar en mini mode por defecto ───
def _iniciar_mini_mode():
    """Activa mini mode inmediatamente."""
    try:
        ventana.update_idletasks()
        abrir_mini_mode()
    except Exception as e:
        _log_error("iniciar_mini_mode", e)
        ventana.deiconify()

ventana.after(500, _iniciar_mini_mode)
ventana.after(200, animar_bola)
ventana.after(30, _procesar_cola_ui)
ventana.mainloop()
